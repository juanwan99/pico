"""Inspect reports irregular school tables without guessing a tidy grid."""

from __future__ import annotations

import io
import sys
from pathlib import Path

import pytest
from openpyxl import Workbook

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "services" / "orchestrator"))

from pico_orchestrator.office.inspect import inspect_office_bytes


def _xlsx_bytes(book: Workbook) -> bytes:
    buf = io.BytesIO()
    book.save(buf)
    return buf.getvalue()


def test_xlsx_two_row_merged_headers_and_window() -> None:
    book = Workbook()
    sheet = book.active
    sheet.title = "全部数据"
    sheet.merge_cells("A1:C1")
    sheet.merge_cells("D1:E1")
    sheet["A1"] = "学生"
    sheet["D1"] = "监护人"
    sheet["A2"] = "学号"
    sheet["B2"] = "姓名"
    sheet["C2"] = "班级"
    sheet["D2"] = "姓名"
    sheet["E2"] = "电话"
    for i in range(3, 23):
        sheet[f"A{i}"] = f"S{i:03d}"
        sheet[f"B{i}"] = f"生{i}"
        sheet[f"C{i}"] = "高一1"
        sheet[f"D{i}"] = f"家{i}"
        sheet[f"E{i}"] = "13800000000"
    raw = _xlsx_bytes(book)
    outline = inspect_office_bytes(
        raw, ".xlsx", header_rows=2, max_rows=5, sheet="全部数据"
    )
    unit = outline["units"][0]
    assert unit["irregular"] is True
    assert "A1:C1" in unit["merges"]
    assert "学生 / 学号" in unit["headers"]
    assert "监护人 / 电话" in unit["headers"]
    assert len(unit["headers"]) == 5
    assert unit["window"]["truncated"] is True
    assert unit["window"]["leftover_rows"] == 15
    assert unit["window"]["leftover_cols"] == 0
    assert unit["window"]["start_col"] == 1
    assert len(unit["preview"]) == 5
    assert unit["preview"][0][0] == "S003"


def test_xlsx_column_window_reports_leftover_cols() -> None:
    book = Workbook()
    sheet = book.active
    sheet.title = "宽表"
    for col in range(1, 11):
        sheet.cell(1, col, f"C{col}")
        sheet.cell(2, col, col)
    raw = _xlsx_bytes(book)
    first = inspect_office_bytes(raw, ".xlsx", max_cols=4, max_rows=8)
    unit = first["units"][0]
    assert unit["cols"] == 10
    assert unit["window"]["start_col"] == 1
    assert unit["window"]["end_col"] == 4
    assert unit["window"]["leftover_cols"] == 6
    assert unit["window"]["truncated"] is True
    assert unit["headers"] == ["C1", "C2", "C3", "C4"]
    next_win = inspect_office_bytes(
        raw, ".xlsx", max_cols=4, start_col=5, max_rows=8
    )
    nxt = next_win["units"][0]
    assert nxt["window"]["start_col"] == 5
    assert nxt["window"]["end_col"] == 8
    assert nxt["window"]["leftover_cols"] == 2
    assert nxt["headers"][0] == "C5"
    last = inspect_office_bytes(
        raw, ".xlsx", max_cols=4, start_col=9, max_rows=8
    )
    done = last["units"][0]
    assert done["window"]["leftover_cols"] == 0
    assert done["window"]["truncated"] is False


def test_xlsx_sheet_filter_and_blank_header_not_guessed() -> None:
    book = Workbook()
    one = book.active
    one.title = "导入说明"
    one["A1"] = "说明"
    one["A2"] = "先看示例再贴数据"
    two = book.create_sheet("教师")
    two["A1"] = "姓名"
    two["B1"] = ""
    two["C1"] = "科目"
    two["A2"] = "李老师"
    two["C2"] = "生物"
    raw = _xlsx_bytes(book)
    outline = inspect_office_bytes(raw, ".xlsx", sheet=2)
    assert outline["sheets"] == 2
    assert len(outline["units"]) == 1
    unit = outline["units"][0]
    assert unit["name"] == "教师"
    assert unit["blank_header_cols"] >= 1
    assert unit["headers"][1] == ""
    with pytest.raises(ValueError, match="工作表"):
        inspect_office_bytes(raw, ".xlsx", sheet="不存在")


def test_docx_nested_table_and_fill_blank_prose() -> None:
    from docx import Document

    with_table = Document()
    with_table.add_paragraph("实验记录")
    table = with_table.add_table(rows=3, cols=3)
    table.cell(0, 0).text = "组别"
    table.cell(0, 1).text = "处理"
    table.cell(0, 2).text = "结果"
    table.cell(1, 0).text = "甲"
    table.cell(1, 1).text = "光照"
    table.cell(1, 2).text = "绿"
    buf = io.BytesIO()
    with_table.save(buf)
    table_outline = inspect_office_bytes(buf.getvalue(), ".docx")
    assert table_outline["layout"] == "tables"
    unit = next(u for u in table_outline["units"] if u["kind"] == "table")
    assert "组别" in unit["headers"]
    assert any("甲" in row for row in unit["preview"])

    prose = Document()
    prose.add_paragraph("一、填空")
    prose.add_paragraph("豌豆的______是隐性性状。")
    prose.add_paragraph("F2 分离比约为______。")
    pbuf = io.BytesIO()
    prose.save(pbuf)
    prose_outline = inspect_office_bytes(pbuf.getvalue(), ".docx")
    assert prose_outline["tables"] == 0
    assert prose_outline["layout"] == "prose"
    texts = [u["text"] for u in prose_outline["units"] if u["kind"] == "para"]
    assert any("填空" in t or "隐性" in t for t in texts)


def test_pptx_table_is_not_screenshot() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "分组"
    table_shape = slide.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(2))
    grid = table_shape.table
    grid.cell(0, 0).text = "组"
    grid.cell(0, 1).text = "人"
    grid.cell(1, 0).text = "甲"
    grid.cell(1, 1).text = "李老师"
    buf = io.BytesIO()
    deck.save(buf)
    outline = inspect_office_bytes(buf.getvalue(), ".pptx")
    assert outline["tables"] == 1
    slide_unit = outline["units"][0]
    assert slide_unit["tables"] == 1
    preview = slide_unit["table_preview"][0]
    assert "组" in preview["headers"]
    assert any("甲" in row for row in preview["preview"])


def test_header_rows_bounds() -> None:
    book = Workbook()
    book.active["A1"] = "x"
    raw = _xlsx_bytes(book)
    with pytest.raises(ValueError, match="header_rows"):
        inspect_office_bytes(raw, ".xlsx", header_rows=0)
    with pytest.raises(ValueError, match="max_rows"):
        inspect_office_bytes(raw, ".xlsx", max_rows=99)


@pytest.mark.asyncio
async def test_inspect_document_passes_header_window() -> None:
    from dataclasses import dataclass
    from typing import Any

    from pico_orchestrator.tools_builtin import build_default_gateway

    @dataclass
    class P:
        school_id: str = "school-a"
        membership_id: str = "member-a"
        scopes: list[str] | None = None

        def __post_init__(self) -> None:
            if self.scopes is None:
                self.scopes = ["ai:run"]

    class Store:
        def __init__(self, raw: bytes) -> None:
            self.row = {
                "artifact_id": "art-1",
                "title": "名册.xlsx",
                "content": raw,
                "kind": "xlsx",
            }

        async def write(self, principal: P, **kwargs: Any) -> dict[str, Any]:
            raise AssertionError("inspect must not write")

        async def read(self, principal: P, *, artifact_id: str | None, title: str | None):
            return self.row

        async def list(self, principal: P, *, limit: int) -> list[dict[str, Any]]:
            return [self.row]

    book = Workbook()
    sheet = book.active
    sheet.title = "全部数据"
    sheet.merge_cells("A1:B1")
    sheet["A1"] = "学生"
    sheet["A2"] = "学号"
    sheet["B2"] = "姓名"
    sheet["A3"] = "S001"
    sheet["B3"] = "甲"
    raw = _xlsx_bytes(book)
    gw = build_default_gateway(Store(raw))
    outline = await gw.invoke(
        P(),
        "inspect_document",
        {
            "artifact_id": "art-1",
            "kind": "xlsx",
            "header_rows": 2,
            "max_rows": 3,
            "start_col": 1,
            "sheet": "全部数据",
        },
    )
    unit = outline["units"][0]
    assert "学生 / 学号" in unit["headers"]
    assert unit["window"]["start_row"] == 3
    assert unit["window"]["leftover_cols"] == 0
