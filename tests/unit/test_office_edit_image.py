"""T-AGENT-EXT-V1: python-docx/pptx edit originals + Zhipu glm-image mock."""

from __future__ import annotations

import base64
import io
import sys
import zipfile
from pathlib import Path
from types import SimpleNamespace

import httpx
import pytest

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "services" / "orchestrator"))

from pico_orchestrator.artifact_types import is_valid_ooxml_package
from pico_orchestrator.document_generators import build_docx_document
from pico_orchestrator.gateway import ToolError
from pico_orchestrator.image_generate import (
    _429_MAX_TRIES,
    _429_RETRY_AFTER_CAP_S,
    IMAGE_TIMEOUT_S,
    NO_KEY_MESSAGE,
    REJECTED_PROVIDER_MESSAGE,
    generate_image_bytes,
    reset_image_generate_runtime,
)
from pico_orchestrator.office_editors import edit_docx_bytes, edit_pptx_title_bytes
from pico_orchestrator.skill_policy import snapshot_for_skill
from pico_orchestrator.tools_builtin import build_default_gateway, openai_tool_schemas
from pico_orchestrator.true_pi.config import ALLOWED_GATEWAY_TOOLS

ONE_PNG = (
    b"\x89PNG\r\n\x1a\n"
    b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
    b"\x00\x00\x00\x0cIDATx\x9cc``\x00\x00\x00\x04\x00\x01\xa3\x0a\x0d\xe4"
    b"\x00\x00\x00\x00IEND\xaeB`\x82"
)


@pytest.fixture(autouse=True)
def _reset_image_runtime(monkeypatch: pytest.MonkeyPatch):
    reset_image_generate_runtime()
    monkeypatch.delenv("GEMINI_API_KEY", raising=False)
    monkeypatch.delenv("GOOGLE_API_KEY", raising=False)
    monkeypatch.delenv("GEMINI_IMAGES_URL", raising=False)
    monkeypatch.delenv("GEMINI_IMAGE_MODEL", raising=False)
    monkeypatch.delenv("PICO_IMAGE_GATEWAY_URL", raising=False)
    monkeypatch.delenv("PICO_IMAGE_GATEWAY_KEY", raising=False)
    monkeypatch.delenv("PICO_IMAGE_GATEWAY_MODEL", raising=False)
    yield
    reset_image_generate_runtime()


def _three_para_docx() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph("第一段保持")
    doc.add_paragraph("第二段也在")
    doc.add_paragraph("第三段很长很长很长很长很长很长很长很长")
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _two_slide_pptx() -> bytes:
    from pptx import Presentation

    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[0])
    first.shapes.title.text = "原标题"
    second = deck.slides.add_slide(deck.slide_layouts[1])
    second.shapes.title.text = "第二页还在"
    out = io.BytesIO()
    deck.save(out)
    return out.getvalue()


def test_edit_docx_keeps_other_paragraphs() -> None:
    raw = _three_para_docx()
    assert is_valid_ooxml_package(raw, ".docx")
    edited = edit_docx_bytes(raw, paragraph_index=3, text="第三段短")
    assert is_valid_ooxml_package(edited, ".docx")
    from docx import Document

    doc = Document(io.BytesIO(edited))
    texts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    assert texts[0] == "第一段保持"
    assert texts[1] == "第二段也在"
    assert texts[2] == "第三段短"
    assert "很长很长" not in texts[2]
    with zipfile.ZipFile(io.BytesIO(edited)) as zf:
        assert "word/document.xml" in zf.namelist()
        assert "word/styles.xml" in zf.namelist()
    generated = build_docx_document(title="x.docx", marker="MARK", body="另造")
    with zipfile.ZipFile(io.BytesIO(generated)) as zf:
        generated_names = set(zf.namelist())
    with zipfile.ZipFile(io.BytesIO(edited)) as zf:
        edited_names = set(zf.namelist())
    assert "word/styles.xml" in edited_names
    assert "word/styles.xml" in generated_names
    assert "MARK" in zf_document_xml(generated)


def zf_document_xml(raw: bytes) -> str:
    with zipfile.ZipFile(io.BytesIO(raw)) as zf:
        return zf.read("word/document.xml").decode("utf-8")


def test_edit_pptx_keeps_other_slides() -> None:
    raw = _two_slide_pptx()
    edited = edit_pptx_title_bytes(raw, slide_index=1, new_title="课堂导入")
    assert is_valid_ooxml_package(edited, ".pptx")
    from pptx import Presentation

    deck = Presentation(io.BytesIO(edited))
    assert deck.slides[0].shapes.title.text.strip() == "课堂导入"
    assert deck.slides[1].shapes.title.text.strip() == "第二页还在"
    assert len(deck.slides) == 2


@pytest.mark.asyncio
async def test_generate_image_no_key_chinese(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    monkeypatch.delenv("ZHIPU_API_KEY", raising=False)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("分数的初步认识")
    assert caught.value.code == "image.unconfigured"
    assert "不能编造" in caught.value.message
    assert "SILICONFLOW" not in caught.value.message
    assert caught.value.message == NO_KEY_MESSAGE


@pytest.mark.asyncio
async def test_generate_image_siliconflow_only_rejected(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("SILICONFLOW_API_KEY", "test-key-not-a-secret")
    monkeypatch.delenv("ZHIPU_API_KEY", raising=False)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("分数的初步认识课堂示意图")
    assert caught.value.code == "image.provider_rejected"
    assert caught.value.message == REJECTED_PROVIDER_MESSAGE
    assert "SILICONFLOW_API_KEY" not in caught.value.message


@pytest.mark.asyncio
async def test_generate_image_zhipu_mock_https_png(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")

    async def fake_post(payload, *, api_key, timeout):
        assert api_key == "test-zhipu-not-a-secret"
        assert payload["model"] == "glm-image"
        assert payload["prompt"]
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "data": [{"b64_json": base64.b64encode(ONE_PNG).decode("ascii")}]
            },
        )

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    raw, ext = await generate_image_bytes("分数的初步认识课堂示意图")
    assert ext == "png"
    assert raw.startswith(b"\x89PNG")


async def _record_sleep(monkeypatch: pytest.MonkeyPatch) -> list[float]:
    reset_image_generate_runtime()
    slept: list[float] = []

    async def fake_sleep(seconds: float) -> None:
        slept.append(float(seconds))

    monkeypatch.setattr("pico_orchestrator.image_generate.asyncio.sleep", fake_sleep)
    return slept


@pytest.mark.asyncio
async def test_generate_image_zhipu_4xx_no_fake(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    slept = await _record_sleep(monkeypatch)
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        return SimpleNamespace(status_code=400, json=lambda: {"error": "bad"})

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("画一只猫")
    assert caught.value.code == "image.provider"
    assert "不能编造" in caught.value.message
    assert calls["n"] == 1
    assert slept == []


@pytest.mark.asyncio
async def test_generate_image_provider_retries_once_then_ok(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Live F2: image.provider on first POST must retry once inside the tool."""
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        if calls["n"] == 1:
            return SimpleNamespace(status_code=500, json=lambda: {"error": "busy"})
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "data": [{"b64_json": base64.b64encode(ONE_PNG).decode("ascii")}]
            },
        )

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    raw, ext = await generate_image_bytes("封面示意图")
    assert calls["n"] == 2
    assert ext == "png"
    assert raw.startswith(b"\x89PNG")


def _429(retry_after: str | None = None, *, code: str | None = None) -> SimpleNamespace:
    headers: dict[str, str] = {}
    if retry_after is not None:
        headers["Retry-After"] = retry_after
    error: dict[str, str] = {"message": "rate"}
    if code is not None:
        error["code"] = code
    return SimpleNamespace(
        status_code=429,
        headers=headers,
        json=lambda: {"error": error},
    )


@pytest.mark.asyncio
async def test_generate_image_429_then_200(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """HTTP 429 with Retry-After rests that long, then the next POST may succeed."""
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    slept = await _record_sleep(monkeypatch)
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        if calls["n"] == 1:
            return _429("4")
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "data": [{"b64_json": base64.b64encode(ONE_PNG).decode("ascii")}]
            },
        )

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    raw, ext = await generate_image_bytes("封面示意图")
    assert calls["n"] == 2
    assert slept == [4.0]
    assert all(s >= 1.0 for s in slept)
    assert ext == "png"
    assert raw.startswith(b"\x89PNG")


@pytest.mark.asyncio
async def test_generate_image_429_exhausted_no_fake(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Short Retry-After still image.provider when every POST 429s; never invent pixels."""
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    slept = await _record_sleep(monkeypatch)
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        return _429("3")

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("封面示意图")
    assert caught.value.code == "image.provider"
    assert "不能编造" in caught.value.message
    assert calls["n"] == _429_MAX_TRIES
    assert slept == [3.0] * (_429_MAX_TRIES - 1)
    assert all(s >= 1.0 for s in slept)


@pytest.mark.asyncio
async def test_generate_image_same_prompt_single_flight(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Live F4: Pi retrying the same cover must not each burn 6 POSTs."""
    import asyncio

    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    reset_image_generate_runtime()
    started = asyncio.Event()
    release = asyncio.Event()
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        started.set()
        await release.wait()
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "data": [{"b64_json": base64.b64encode(ONE_PNG).decode("ascii")}]
            },
        )

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    first = asyncio.create_task(generate_image_bytes("corporate cover"))
    await started.wait()
    second = asyncio.create_task(generate_image_bytes("corporate cover"))
    await asyncio.sleep(0)
    release.set()
    one, two = await asyncio.gather(first, second)
    assert calls["n"] == 1
    assert one == two
    assert one[1] == "png"


@pytest.mark.asyncio
async def test_generate_image_429_exhausted_next_call_uses_gate(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """After 429 with no Retry-After, the next POST waits the timeout window."""
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    slept = await _record_sleep(monkeypatch)
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        if calls["n"] == 1:
            return _429()
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "data": [{"b64_json": base64.b64encode(ONE_PNG).decode("ascii")}]
            },
        )

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("封面示意图")
    assert caught.value.code == "image.provider"
    assert "不能编造" in caught.value.message
    assert calls["n"] == 1
    assert slept == []
    raw, ext = await generate_image_bytes("另一张封面")
    assert ext == "png"
    assert raw.startswith(b"\x89PNG")
    assert calls["n"] == 2
    assert slept == [pytest.approx(_429_RETRY_AFTER_CAP_S, abs=0.05)]


@pytest.mark.asyncio
async def test_generate_image_429_missing_header_does_not_consecutive_fly(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Live F5: no Retry-After must not 2/4/8/16/30 then ~28s then fly again."""
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    slept = await _record_sleep(monkeypatch)
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        return _429()

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("封面示意图")
    assert caught.value.code == "image.provider"
    with pytest.raises(ToolError) as caught2:
        await generate_image_bytes("另一张封面")
    assert caught2.value.code == "image.provider"
    assert "不能编造" in caught2.value.message
    assert calls["n"] == 2
    assert slept == [pytest.approx(IMAGE_TIMEOUT_S, abs=0.05)]


@pytest.mark.asyncio
async def test_generate_image_429_retry_after_capped_at_timeout_window(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    slept = await _record_sleep(monkeypatch)
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        if calls["n"] == 1:
            return _429("120")
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "data": [{"b64_json": base64.b64encode(ONE_PNG).decode("ascii")}]
            },
        )

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("封面示意图")
    assert caught.value.code == "image.provider"
    assert calls["n"] == 1
    _raw, ext = await generate_image_bytes("另一张封面")
    assert ext == "png"
    assert calls["n"] == 2
    assert slept == [pytest.approx(IMAGE_TIMEOUT_S, abs=0.05)]


@pytest.mark.asyncio
async def test_generate_image_429_retry_after_http_date(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from datetime import UTC, datetime, timedelta
    from email.utils import format_datetime

    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    slept = await _record_sleep(monkeypatch)
    calls = {"n": 0}
    when = format_datetime(datetime.now(UTC) + timedelta(seconds=8))

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        if calls["n"] == 1:
            return _429(when)
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "data": [{"b64_json": base64.b64encode(ONE_PNG).decode("ascii")}]
            },
        )

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    _raw, ext = await generate_image_bytes("封面示意图")
    assert calls["n"] == 2
    assert ext == "png"
    assert len(slept) == 1
    assert 6.0 <= slept[0] <= 8.5


@pytest.mark.asyncio
async def test_generate_image_1113_no_retry_no_fake(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Live F5 body: error.code=1113 余额不足 — not a rate limit. One POST."""
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    slept = await _record_sleep(monkeypatch)
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        return _429("4", code="1113")

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("封面示意图")
    assert caught.value.code == "image.provider"
    assert "不能编造" in caught.value.message
    assert calls["n"] == 1
    assert slept == []
    with pytest.raises(ToolError):
        await generate_image_bytes("另一张封面")
    assert calls["n"] == 2
    assert slept == [pytest.approx(IMAGE_TIMEOUT_S, abs=0.05)]


@pytest.mark.asyncio
async def test_generate_image_gemini_mock_https_png(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("GEMINI_API_KEY", "test-gemini-not-a-secret")
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    seen: dict[str, object] = {}

    async def fake_gemini(payload, *, api_key, timeout):
        seen["key"] = api_key
        seen["modalities"] = payload["generationConfig"]["responseModalities"]
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "candidates": [
                    {
                        "content": {
                            "parts": [
                                {
                                    "inlineData": {
                                        "mimeType": "image/png",
                                        "data": base64.b64encode(ONE_PNG).decode(
                                            "ascii"
                                        ),
                                    }
                                }
                            ]
                        }
                    }
                ]
            },
        )

    async def zhipu_should_not_run(payload, *, api_key, timeout):
        raise AssertionError("zhipu must not run when Gemini key is set")

    monkeypatch.setattr("pico_orchestrator.image_generate._post_gemini", fake_gemini)
    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", zhipu_should_not_run)
    raw, ext = await generate_image_bytes("封面示意图")
    assert seen["key"] == "test-gemini-not-a-secret"
    assert "IMAGE" in seen["modalities"]
    assert ext == "png"
    assert raw.startswith(b"\x89PNG")


@pytest.mark.asyncio
async def test_generate_image_gateway_mock_https_png(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("PICO_IMAGE_GATEWAY_URL", "https://newapi.example.com")
    monkeypatch.setenv("PICO_IMAGE_GATEWAY_KEY", "sk-gateway-not-a-secret")
    monkeypatch.setenv("PICO_IMAGE_GATEWAY_MODEL", "gemini-2.5-flash-image")
    monkeypatch.setenv("GEMINI_API_KEY", "test-gemini-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    seen: dict[str, object] = {}

    async def fake_gateway(payload, *, api_key, timeout, url=None):
        seen["key"] = api_key
        seen["url"] = url
        seen["payload"] = payload
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "candidates": [
                    {
                        "content": {
                            "parts": [
                                {
                                    "inlineData": {
                                        "mimeType": "image/png",
                                        "data": base64.b64encode(ONE_PNG).decode(
                                            "ascii"
                                        ),
                                    }
                                }
                            ]
                        }
                    }
                ]
            },
        )

    async def gemini_should_not_run(payload, *, api_key, timeout):
        raise AssertionError("direct Gemini must not run when gateway is set")

    monkeypatch.setattr("pico_orchestrator.image_generate._post_gateway", fake_gateway)
    monkeypatch.setattr(
        "pico_orchestrator.image_generate._post_gemini", gemini_should_not_run
    )
    raw, ext = await generate_image_bytes("封面示意图")
    assert seen["key"] == "sk-gateway-not-a-secret"
    assert str(seen["url"]).endswith(
        "/v1beta/models/gemini-2.5-flash-image:generateContent"
    )
    payload = seen["payload"]
    assert isinstance(payload, dict)
    assert payload["contents"][0]["parts"][0]["text"] == "封面示意图"
    assert "IMAGE" in payload["generationConfig"]["responseModalities"]
    assert ext == "png"
    assert raw.startswith(b"\x89PNG")


@pytest.mark.asyncio
async def test_generate_image_gateway_imagen_openai_path(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("PICO_IMAGE_GATEWAY_URL", "https://newapi.example.com")
    monkeypatch.setenv("PICO_IMAGE_GATEWAY_KEY", "sk-gateway-not-a-secret")
    monkeypatch.setenv("PICO_IMAGE_GATEWAY_MODEL", "imagen-4.0-generate-001")
    monkeypatch.delenv("GEMINI_API_KEY", raising=False)
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    seen: dict[str, object] = {}

    async def fake_gateway(payload, *, api_key, timeout, url=None):
        seen["url"] = url
        seen["payload"] = payload
        return SimpleNamespace(
            status_code=200,
            json=lambda: {
                "data": [{"b64_json": base64.b64encode(ONE_PNG).decode("ascii")}]
            },
        )

    monkeypatch.setattr("pico_orchestrator.image_generate._post_gateway", fake_gateway)
    raw, ext = await generate_image_bytes("封面示意图")
    assert seen["url"] is None
    payload = seen["payload"]
    assert isinstance(payload, dict)
    assert payload["model"] == "imagen-4.0-generate-001"
    assert payload["prompt"] == "封面示意图"
    assert ext == "png"
    assert raw.startswith(b"\x89PNG")


@pytest.mark.asyncio
async def test_generate_image_gemini_400_no_fake(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("GEMINI_API_KEY", "test-gemini-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)

    async def fake_gemini(payload, *, api_key, timeout):
        return SimpleNamespace(status_code=400, json=lambda: {"error": {"message": "bad"}})

    monkeypatch.setattr("pico_orchestrator.image_generate._post_gemini", fake_gemini)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("封面示意图")
    assert caught.value.code == "image.provider"
    assert "不能编造" in caught.value.message


@pytest.mark.asyncio
async def test_generate_image_timeout_no_retry(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("ZHIPU_API_KEY", "test-zhipu-not-a-secret")
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    slept = await _record_sleep(monkeypatch)
    calls = {"n": 0}

    async def fake_post(payload, *, api_key, timeout):
        calls["n"] += 1
        raise httpx.TimeoutException("read timeout")

    monkeypatch.setattr("pico_orchestrator.image_generate._post_images", fake_post)
    with pytest.raises(ToolError) as caught:
        await generate_image_bytes("封面示意图")
    assert caught.value.code == "image.timeout"
    assert "不能编造" in caught.value.message
    assert calls["n"] == 1
    assert slept == []


def test_sidebar_chat_has_no_edit_or_image_tools() -> None:
    chat = snapshot_for_skill("skill.chat")
    assert chat is not None
    assert chat["tools"] == []
    assert openai_tool_schemas(build_default_gateway(), allowed_tools=[]) == []
    forbidden = {
        "edit_docx_document",
        "edit_pptx_document",
        "generate_image",
        "generate_diagram",
    }
    assert forbidden.isdisjoint(chat["tools"])
    deliver = snapshot_for_skill("skill-deliverable")
    assert deliver is not None
    assert forbidden <= set(deliver["tools"])
    assert forbidden <= ALLOWED_GATEWAY_TOOLS
    gw_names = set(build_default_gateway().tools)
    assert forbidden <= gw_names
    ts = (ROOT / "services" / "true_pi_bridge" / "pico-gateway-tools.ts").read_text(
        encoding="utf-8"
    )
    for name in forbidden:
        assert f'"{name}"' in ts
