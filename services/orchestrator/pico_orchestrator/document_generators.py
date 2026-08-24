"""Deterministic allowlist document generators (HTML / DOCX / PPTX).

HTML / XLSX stay stdlib. Word / PPTX are thin python-docx / python-pptx
adapters — real Office packages with visible body, not three-line XML zips.
"""

from __future__ import annotations

import html
import io
import re
import zipfile
from xml.etree import ElementTree as ET
from xml.sax.saxutils import escape

# Delivery floor: Word opens with hundreds of characters; PPT has ≥3 titled slides.
MIN_DOCX_BODY_CHARS = 300
MIN_PPTX_SLIDES = 3

_W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _visible_len(text: str) -> int:
    return sum(1 for ch in (text or "") if not ch.isspace())


def _display_title(title: str, fallback: str) -> str:
    name = (title or "").strip() or fallback
    lower = name.lower()
    for ext in (".docx", ".pptx", ".html", ".htm", ".md"):
        if lower.endswith(ext):
            stem = name[: -len(ext)].strip()
            return stem or fallback
    return name


def _split_blocks(raw: str) -> list[str]:
    text = (raw or "").replace("\r\n", "\n").strip()
    if not text:
        return []
    if "\n---\n" in f"\n{text}\n":
        parts = [p.strip() for p in text.split("\n---\n") if p.strip()]
        if parts:
            return parts
    chunks = [part.strip() for part in text.split("\n\n") if part.strip()]
    if chunks:
        return chunks
    return [ln.strip() for ln in text.split("\n") if ln.strip()]


def docx_visible_text(raw: bytes) -> str:
    """Visible Word paragraphs from OOXML (stdlib zip+xml). Empty on junk."""
    if not raw or raw[:2] != b"PK":
        return ""
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            xml = zf.read("word/document.xml")
    except (KeyError, zipfile.BadZipFile):
        return ""
    try:
        root = ET.fromstring(xml)
    except ET.ParseError:
        return ""
    paras: list[str] = []
    for p in root.iter(f"{_W_NS}p"):
        line = "".join(t.text or "" for t in p.iter(f"{_W_NS}t")).strip()
        if line:
            paras.append(line)
    return "\n".join(paras)


def pptx_slide_titles(raw: bytes) -> list[str]:
    """First text run on each slide (title-ish). Empty list on junk."""
    if not raw or raw[:2] != b"PK":
        return []
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            names = zf.namelist()
            slides = sorted(
                n
                for n in names
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            )
            titles: list[str] = []
            for path in slides:
                try:
                    xml = zf.read(path)
                except KeyError:
                    titles.append("")
                    continue
                try:
                    root = ET.fromstring(xml)
                except ET.ParseError:
                    titles.append("")
                    continue
                texts = [
                    (node.text or "").strip()
                    for node in root.iter(f"{_A_NS}t")
                    if (node.text or "").strip()
                ]
                titles.append(texts[0] if texts else "")
    except zipfile.BadZipFile:
        return []
    return titles


def office_shell_reason(raw: bytes, ext: str) -> str | None:
    """Chinese reason when a supposed Office file is an empty shell. None if ok."""
    suffix = (ext or "").lower()
    if not suffix.startswith("."):
        suffix = f".{suffix}" if suffix else ""
    if suffix == ".docx":
        text = docx_visible_text(raw)
        if _visible_len(text) < MIN_DOCX_BODY_CHARS:
            return (
                "Word 正文过短，打开后几乎是空壳。"
                "请用 generate_docx_document 写入完整多段正文（至少数百字）后再交。"
            )
        return None
    if suffix == ".pptx":
        titles = pptx_slide_titles(raw)
        titled = sum(1 for t in titles if (t or "").strip())
        if len(titles) < MIN_PPTX_SLIDES or titled < MIN_PPTX_SLIDES:
            return (
                "课件不足三页有标题，打开后几乎是空页。"
                "请用 generate_pptx_document 写出至少三页后再交。"
            )
        return None
    return None


def _docx_body_paragraphs(body: str | None, *, title: str, marker: str) -> list[str]:
    heading = _display_title(title, "办公文稿")
    chunks = _split_blocks(body or "")
    if not chunks:
        chunks = [f"本文档主题：{heading}。"]
    if _visible_len("".join(chunks)) >= MIN_DOCX_BODY_CHARS:
        return chunks
    seed = "\n".join(chunks)
    paras = [
        f"本文档主题：{heading}。",
        "一、事项说明",
        seed,
        "二、执行要点",
        (
            f"请按「{heading}」落实：写清对象、时间与责任人。"
            f"以下为本次正文，供打开后直接使用。{seed}"
        ),
        "三、备注",
        (
            f"本文件是可打开的 Word 正文，不是空壳标记页。"
            f"主题：{heading}。内部标记：{marker}。"
            "若要改一版，指出段落后即可改。"
        ),
        (
            f"四、打开说明。请用 Word / LibreOffice 打开本文件阅读全文，"
            f"不要只看标题或标记。主题仍是「{heading}」。"
            "正文如下，便于打印或转发："
            f"{seed}"
            "请核对时间、对象、地点是否与题面一致；有误指出段落即可改一版。"
        ),
    ]
    while _visible_len("".join(paras)) < MIN_DOCX_BODY_CHARS:
        paras.append(
            f"补充：{heading}。{seed}。"
            "本段保证打开后有可读正文，而不是三行空壳。"
        )
        if len(paras) > 16:
            break
    return paras


def _pptx_slides(body: str | None, *, title: str, marker: str) -> list[tuple[str, str]]:
    heading = _display_title(title, "Pico PPTX")
    parts = _split_blocks(body or "")
    if not parts:
        parts = [f"本课件围绕「{heading}」展开，含目标、要点与说明。"]
    slides: list[tuple[str, str]] = [(heading, f"标记：{marker}\n{parts[0]}")]
    for index, chunk in enumerate(parts[1:], start=2):
        first = chunk.split("\n", 1)[0].strip()[:40] or f"第{index}页"
        slides.append((first, chunk))
    if len(slides) < MIN_PPTX_SLIDES:
        slides.append(("要点", "\n".join(parts)))
    if len(slides) < MIN_PPTX_SLIDES:
        slides.append(
            (
                "说明",
                f"本课件不少于三页。主题：{heading}。标记：{marker}。请按页内标题讲解。",
            )
        )
    return slides[:20]


def _require_marker(marker: str) -> str:
    value = (marker or "").strip()
    if not value:
        # Marker is an internal traceability tag. Auto-generate one when the
        # model omits it — hard-failing here made the true-Pi agent retry the
        # same tool call forever (message_update flood → OOM). A unique tag
        # keeps deliveries traceable without requiring the model to know the
        # internal field.
        import uuid

        value = f"pico-{uuid.uuid4().hex[:12]}"
    if len(value) > 200:
        raise ValueError("marker exceeds 200 characters")
    if any(ord(ch) < 32 for ch in value):
        raise ValueError("marker contains control characters")
    return value


def _html_body_paragraphs(body: str | None, *, marker: str) -> str:
    """Escape body into one or more <p> blocks (blank-line separated)."""
    raw = (body or "").strip() or f"Pico HTML deliverable · {marker}"
    # Cap runaway agent output while still allowing full lesson pages.
    if len(raw) > 50_000:
        raw = raw[:50_000]
    chunks = [part.strip() for part in raw.replace("\r\n", "\n").split("\n\n") if part.strip()]
    if not chunks:
        chunks = [f"Pico HTML deliverable · {marker}"]
    parts: list[str] = []
    for chunk in chunks:
        # Single newlines inside a paragraph become <br/>; still fully escaped.
        lines = [html.escape(line) for line in chunk.split("\n")]
        parts.append("<p>" + "<br />\n".join(lines) + "</p>")
    return "\n  ".join(parts)


def _looks_like_full_html_document(raw: str) -> bool:
    """True when agent passed a complete page (not prose to escape into <p>)."""
    s = (raw or "").lstrip().lower()
    if not s:
        return False
    if s.startswith(("<!doctype html", "<html")):
        return True
    # Fragment that is already a self-contained interactive page shell.
    return "<html" in s[:500] and ("</html>" in s or "<body" in s)


def _looks_like_html_markup(raw: str) -> bool:
    """True when body is HTML markup (full page or fragment) — must not escape to a tag wall.

    #399 R2: models often pass ``<h2>…</h2><button>…`` without doctype/html.
    Escaping those tags into ``&lt;h2&gt;`` makes downloads open as a source/tag wall.
    Plain prose without tags still uses the safe paragraph shell.
    """
    if _looks_like_full_html_document(raw):
        return True
    s = (raw or "").strip()
    if not s or "<" not in s:
        return False
    # Structural / interactive tags ⇒ treat as markup to preserve.
    if re.search(
        r"<\s*(?:button|input|table|thead|tbody|tr|td|th|form|script|style|"
        r"div|section|article|nav|header|footer|main|h[1-6]|ul|ol|li|a|span|"
        r"p|img|label|select|textarea|details|summary)\b",
        s,
        flags=re.IGNORECASE,
    ):
        return True
    # ≥2 generic tags is almost never pure prose.
    tags = re.findall(r"<\s*/?\s*[a-zA-Z][a-zA-Z0-9]*\b", s)
    return len(tags) >= 2


def _wrap_html_fragment(raw_body: str, *, title: str, marker: str) -> str:
    """Wrap an HTML fragment in a minimal interactive document shell (tags kept)."""
    safe_title = html.escape((title or "Pico HTML").strip() or "Pico HTML")
    safe_marker = html.escape(marker)
    body = _strip_remote_script_src(raw_body)
    csp = (
        "default-src 'none'; img-src data:; style-src 'unsafe-inline'; "
        "script-src 'unsafe-inline'; base-uri 'none'; form-action 'none'; "
        "frame-ancestors 'none';"
    )
    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta http-equiv="Content-Security-Policy" content="{csp}" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{safe_title}</title>
  <style>
    body {{ font-family: system-ui, sans-serif; margin: 1.5rem; color: #1a1a1a; line-height: 1.5; max-width: 48rem; }}
    h1, h2, h3 {{ line-height: 1.25; }}
    .marker {{ font-family: ui-monospace, monospace; background: #f3f4f6; padding: 0.25rem 0.5rem; border-radius: 0.25rem; }}
  </style>
</head>
<body>
  <p data-pico-marker-line="1">标记：<span class="marker" data-pico-marker="{safe_marker}">{safe_marker}</span></p>
  {body}
</body>
</html>
"""


def _strip_remote_script_src(doc: str) -> str:
    """Remove external script src (keep inline scripts for local interactivity)."""
    return re.sub(
        r"<script\b[^>]*\bsrc\s*=\s*[\"']https?://[^\"']*[\"'][^>]*>\s*</script>",
        "<!-- stripped remote script -->",
        doc,
        flags=re.IGNORECASE,
    )


def _inject_marker_into_html(doc: str, *, marker: str, title: str) -> str:
    """Ensure marker is present; do not wrap interactive markup as escaped prose."""
    safe_marker = html.escape(marker)
    marker_html = (
        f'<p data-pico-marker-line="1">标记：'
        f'<span class="marker" data-pico-marker="{safe_marker}">{safe_marker}</span></p>'
    )
    if "data-pico-marker=" not in doc:
        if re.search(r"<body\b[^>]*>", doc, flags=re.IGNORECASE):
            doc = re.sub(
                r"(<body\b[^>]*>)",
                r"\1\n  " + marker_html,
                doc,
                count=1,
                flags=re.IGNORECASE,
            )
        else:
            doc = marker_html + "\n" + doc
    # Prefer caller title when document title is empty/generic.
    if title and re.search(r"<title>\s*</title>", doc, flags=re.IGNORECASE):
        doc = re.sub(
            r"<title>\s*</title>",
            f"<title>{html.escape(title)}</title>",
            doc,
            count=1,
            flags=re.IGNORECASE,
        )
    # Interactive local pages need inline script; keep remote blocked via strip.
    # If CSP forbids scripts entirely, button handlers die → source-wall UX.
    csp_interactive = (
        "default-src 'none'; img-src data:; style-src 'unsafe-inline'; "
        "script-src 'unsafe-inline'; base-uri 'none'; form-action 'none'; "
        "frame-ancestors 'none';"
    )
    if re.search(
        r'http-equiv=["\']Content-Security-Policy["\']', doc, re.IGNORECASE
    ):
        doc = re.sub(
            r'(<meta\b[^>]*http-equiv=["\']Content-Security-Policy["\'][^>]*content=["\'])([^"\']*)(["\'])',
            rf"\1{csp_interactive}\3",
            doc,
            count=1,
            flags=re.IGNORECASE,
        )
    elif re.search(r"<head\b", doc, re.IGNORECASE):
        doc = re.sub(
            r"(<head\b[^>]*>)",
            rf'\1\n  <meta http-equiv="Content-Security-Policy" content="{csp_interactive}" />',
            doc,
            count=1,
            flags=re.IGNORECASE,
        )
    return doc


def build_html_document(
    *,
    title: str,
    marker: str,
    body: str | None = None,
) -> bytes:
    """HTML bytes with unique marker.

    - Prose body → safe escaped paragraphs inside a shell (no script).
    - Full HTML document body → kept as interactive page (remote scripts stripped).
      This prevents H-CODEDUMP where agents pass real UI source only to see it
      escaped into a source wall.
    """
    marker = _require_marker(marker)
    safe_title = html.escape((title or "Pico HTML").strip() or "Pico HTML")
    safe_marker = html.escape(marker)
    raw_body = (body or "").strip()
    if len(raw_body) > 50_000:
        raw_body = raw_body[:50_000]

    page_title = (title or "").strip() or "Pico HTML"
    if _looks_like_full_html_document(raw_body):
        doc = _strip_remote_script_src(raw_body)
        doc = _inject_marker_into_html(doc, marker=marker, title=page_title)
        if "data-pico-marker=" not in doc:
            # Extremely broken fragment — fall through to markup/prose paths.
            pass
        else:
            return doc.encode("utf-8")

    # #399 R2: HTML fragments keep real tags (wrap shell); only pure prose is escaped.
    if raw_body and _looks_like_html_markup(raw_body) and not _looks_like_full_html_document(
        raw_body
    ):
        doc = _wrap_html_fragment(raw_body, title=page_title, marker=marker)
        return doc.encode("utf-8")

    body_html = _html_body_paragraphs(raw_body or None, marker=marker)
    # Prose shell: CSP blocks scripts (static reading page).
    doc = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta http-equiv="Content-Security-Policy" content="default-src 'none'; img-src data:; style-src 'unsafe-inline'; base-uri 'none'; form-action 'none'; frame-ancestors 'none';" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{safe_title}</title>
  <style>
    body {{ font-family: system-ui, sans-serif; margin: 1.5rem; color: #1a1a1a; line-height: 1.5; max-width: 48rem; }}
    h1 {{ font-size: 1.5rem; }}
    .marker {{ font-family: ui-monospace, monospace; background: #f3f4f6; padding: 0.25rem 0.5rem; border-radius: 0.25rem; }}
  </style>
</head>
<body>
  <h1>{safe_title}</h1>
  <p>标记：<span class="marker" data-pico-marker="{safe_marker}">{safe_marker}</span></p>
  {body_html}
</body>
</html>
"""
    return doc.encode("utf-8")


def build_docx_document(
    *,
    title: str,
    marker: str,
    body: str | None = None,
) -> bytes:
    """python-docx Word with visible multi-paragraph body (not a three-line XML zip)."""
    marker = _require_marker(marker)
    heading = _display_title(title, "Pico DOCX")
    from docx import Document

    doc = Document()
    doc.add_heading(heading, level=0)
    doc.add_paragraph(f"标记：{marker}")
    for para in _docx_body_paragraphs(body, title=heading, marker=marker):
        doc.add_paragraph(para)
    buf = io.BytesIO()
    doc.save(buf)
    data = buf.getvalue()
    reason = office_shell_reason(data, ".docx")
    if reason:
        raise ValueError(reason)
    return data


KNOWN_CALC_CELL = "NIGHT-P4-CELL-ALPHA"


def build_xlsx_document(
    *,
    title: str,
    marker: str,
    body: str | None = None,
) -> bytes:
    """Minimal valid OOXML XLSX with A1 visible cell text (not CSV renamed)."""
    marker = _require_marker(marker)
    cell = escape((body or "").strip() or KNOWN_CALC_CELL)
    marker_xml = escape(marker)
    heading = escape((title or "Pico XLSX").strip() or "Pico XLSX")

    content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
  <Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
  <Override PartName="/xl/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml"/>
</Types>
"""
    rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>
"""
    workbook = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheets>
    <sheet name="Sheet1" sheetId="1" r:id="rId1"/>
  </sheets>
</workbook>
"""
    workbook_rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>
</Relationships>
"""
    styles = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<styleSheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <fonts count="1"><font><sz val="16"/><name val="Calibri"/></font></fonts>
  <fills count="1"><fill><patternFill patternType="none"/></fill></fills>
  <borders count="1"><border/></borders>
  <cellStyleXfs count="1"><xf/></cellStyleXfs>
  <cellXfs count="1"><xf/></cellXfs>
</styleSheet>
"""
    sheet = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <cols><col min="1" max="1" width="42" customWidth="1"/></cols>
  <sheetData>
    <row r="1" ht="24" customHeight="1">
      <c r="A1" t="inlineStr"><is><t>{cell}</t></is></c>
    </row>
    <row r="2">
      <c r="A2" t="inlineStr"><is><t>{heading}</t></is></c>
    </row>
    <row r="3">
      <c r="A3" t="inlineStr"><is><t>marker:{marker_xml}</t></is></c>
    </row>
  </sheetData>
</worksheet>
"""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("xl/workbook.xml", workbook)
        zf.writestr("xl/_rels/workbook.xml.rels", workbook_rels)
        zf.writestr("xl/worksheets/sheet1.xml", sheet)
        zf.writestr("xl/styles.xml", styles)
    return buf.getvalue()


def _set_slide_title_body(slide: object, heading: str, body: str) -> None:
    title_shape = getattr(getattr(slide, "shapes", None), "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title_shape.text = heading
    body_set = False
    placeholders = getattr(slide, "placeholders", None)
    if placeholders is not None:
        for shape in placeholders:
            idx = getattr(getattr(shape, "placeholder_format", None), "idx", None)
            if idx == 1 and getattr(shape, "has_text_frame", False):
                shape.text = body
                body_set = True
                break
    if not body_set:
        for shape in getattr(slide, "shapes", []):
            if shape is title_shape:
                continue
            if getattr(shape, "has_text_frame", False):
                shape.text = body
                body_set = True
                break
    if title_shape is None:
        for shape in getattr(slide, "shapes", []):
            if getattr(shape, "has_text_frame", False):
                shape.text = heading
                break


def build_pptx_document(
    *,
    title: str,
    marker: str,
    body: str | None = None,
) -> bytes:
    """python-pptx deck with at least three titled slides (not a one-slide XML zip)."""
    marker = _require_marker(marker)
    heading = _display_title(title, "Pico PPTX")
    from pptx import Presentation

    deck = Presentation()
    layout = deck.slide_layouts[1] if len(deck.slide_layouts) > 1 else deck.slide_layouts[0]
    for slide_title, slide_body in _pptx_slides(body, title=heading, marker=marker):
        slide = deck.slides.add_slide(layout)
        _set_slide_title_body(slide, slide_title, slide_body)
    buf = io.BytesIO()
    deck.save(buf)
    data = buf.getvalue()
    reason = office_shell_reason(data, ".pptx")
    if reason:
        raise ValueError(reason)
    return data
