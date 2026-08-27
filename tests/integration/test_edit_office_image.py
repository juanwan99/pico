"""Workbench edit-original + image download (real OOXML, mocked HTTPS)."""

from __future__ import annotations

import base64
import io
import sys
from pathlib import Path

import pytest
from fastapi.testclient import TestClient

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "services" / "api"))
sys.path.insert(0, str(ROOT / "services" / "orchestrator"))

from app.main import app
from pico_orchestrator.artifact_types import is_valid_ooxml_package

ONE_PNG = (
    b"\x89PNG\r\n\x1a\n"
    b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
    b"\x00\x00\x00\x0cIDATx\x9cc``\x00\x00\x00\x04\x00\x01\xa3\x0a\x0d\xe4"
    b"\x00\x00\x00\x00IEND\xaeB`\x82"
)


@pytest.fixture()
def client(tmp_path, monkeypatch):
    database = tmp_path / "edit-office.db"
    monkeypatch.setenv("PICO_DATABASE_URL", f"sqlite+aiosqlite:///{database}")
    monkeypatch.setenv("PICO_JWT_SECRET", "test-secret-at-least-32-bytes-long!!")
    monkeypatch.setenv("PICO_ENV", "development")
    from app import db as dbmod
    from app.settings import get_settings

    get_settings.cache_clear()
    dbmod._engine = None
    dbmod._Session = None
    with TestClient(app) as test_client:
        yield test_client


def _headers(client: TestClient) -> dict[str, str]:
    response = client.post(
        "/v1/dev/token",
        json={"school_id": "school-a", "membership_id": "member-a"},
    )
    assert response.status_code == 200, response.text
    return {"Authorization": f"Bearer {response.json()['access_token']}"}


def _invoke(client: TestClient, headers: dict[str, str], name: str, arguments: dict):
    return client.post(
        "/v1/tools/invoke",
        headers=headers,
        json={"name": name, "arguments": arguments},
    )


def _three_para_docx() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph("第一段保持")
    doc.add_paragraph("第二段也在")
    doc.add_paragraph("第三段很长很长很长很长很长很长很长很长")
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _two_slide_pptx() -> bytes:
    from pptx import Presentation

    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[0])
    first.shapes.title.text = "原标题"
    second = deck.slides.add_slide(deck.slide_layouts[1])
    second.shapes.title.text = "第二页还在"
    out = io.BytesIO()
    deck.save(out)
    return out.getvalue()


def test_upload_edit_docx_download(client) -> None:
    headers = _headers(client)
    raw = _three_para_docx()
    uploaded = client.post(
        "/v1/files",
        headers=headers,
        json={
            "filename": "三段.docx",
            "content_b64": base64.b64encode(raw).decode("ascii"),
        },
    )
    assert uploaded.status_code == 200, uploaded.text
    src_id = uploaded.json()["id"]
    assert src_id
    edited = _invoke(
        client,
        headers,
        "edit_docx_document",
        {"artifact_id": src_id, "paragraph_index": 3, "text": "第三段短"},
    )
    assert edited.status_code == 200, edited.text
    result = edited.json()["result"]
    aid = result["artifact_id"]
    content = client.get(
        f"/v1/artifacts/{aid}/content?download=true",
        headers=headers,
    )
    assert content.status_code == 200
    blob = content.content
    assert is_valid_ooxml_package(blob, ".docx")
    from docx import Document

    doc = Document(io.BytesIO(blob))
    texts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    assert texts[0] == "第一段保持"
    assert texts[1] == "第二段也在"
    assert texts[2] == "第三段短"


def test_upload_edit_pptx_download(client) -> None:
    headers = _headers(client)
    raw = _two_slide_pptx()
    uploaded = client.post(
        "/v1/files",
        headers=headers,
        json={
            "filename": "两页.pptx",
            "content_b64": base64.b64encode(raw).decode("ascii"),
        },
    )
    assert uploaded.status_code == 200, uploaded.text
    src_id = uploaded.json()["id"]
    edited = _invoke(
        client,
        headers,
        "edit_pptx_document",
        {"artifact_id": src_id, "slide_index": 1, "new_title": "课堂导入"},
    )
    assert edited.status_code == 200, edited.text
    aid = edited.json()["result"]["artifact_id"]
    content = client.get(
        f"/v1/artifacts/{aid}/content?download=true",
        headers=headers,
    )
    blob = content.content
    assert is_valid_ooxml_package(blob, ".pptx")
    from pptx import Presentation

    deck = Presentation(io.BytesIO(blob))
    assert deck.slides[0].shapes.title.text.strip() == "课堂导入"
    assert deck.slides[1].shapes.title.text.strip() == "第二页还在"


def test_generate_image_siliconflow_key_rejected(client, monkeypatch) -> None:
    monkeypatch.setenv("SILICONFLOW_API_KEY", "test-key-not-a-secret")
    monkeypatch.delenv("ZHIPU_API_KEY", raising=False)
    headers = _headers(client)
    created = _invoke(
        client,
        headers,
        "generate_image",
        {"prompt": "分数的初步认识课堂示意图", "title": "分数.png"},
    )
    assert created.status_code == 400, created.text
    detail = created.json().get("detail") or {}
    message = detail.get("message") if isinstance(detail, dict) else str(detail)
    assert "否决" in str(message)
    assert "SILICONFLOW_API_KEY" not in str(message)
    assert "artifact_id" not in created.text


def test_generate_image_no_key_no_artifact(client, monkeypatch) -> None:
    monkeypatch.delenv("SILICONFLOW_API_KEY", raising=False)
    monkeypatch.delenv("ZHIPU_API_KEY", raising=False)
    headers = _headers(client)
    created = _invoke(
        client,
        headers,
        "generate_image",
        {"prompt": "分数的初步认识课堂示意图"},
    )
    assert created.status_code == 400
    detail = created.json().get("detail") or {}
    if isinstance(detail, dict):
        message = detail.get("message") or ""
    else:
        message = str(detail)
    assert "不能编造" in message
    assert "SILICONFLOW" not in message
    assert "artifact_id" not in created.text
