"""Pull embedded picture bytes from uploaded Office files. Path B.

No spec rebuild. No HTTP fetch. Pixels stay pixels.
"""

from __future__ import annotations

import io

from pico_orchestrator.office.legacy import require_supported_office_ext

_MAX_IMAGES = 8
_MAX_IMAGE_BYTES = 8 * 1024 * 1024
_PNG = b"\x89PNG\r\n\x1a\n"
_JPEG = b"\xff\xd8"


def extract_embedded_images(raw: bytes, ext: str) -> list[bytes]:
    """Return embedded png/jpeg blobs. Empty for sheets / no pictures."""
    suffix = require_supported_office_ext(ext)
    if suffix == ".pptx":
        return _from_pptx(raw)
    if suffix == ".docx":
        return _from_docx(raw)
    return []


def _keep(blob: bytes | None) -> bytes | None:
    if not blob or len(blob) > _MAX_IMAGE_BYTES:
        return None
    if blob.startswith(_PNG) or blob[:2] == _JPEG:
        return blob
    return None


def _from_pptx(raw: bytes) -> list[bytes]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = Presentation(io.BytesIO(raw))
    out: list[bytes] = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = _keep(bytes(shape.image.blob))
            except Exception:  # noqa: BLE001 — skip a broken picture part
                blob = None
            if blob is None:
                continue
            out.append(blob)
            if len(out) >= _MAX_IMAGES:
                return out
    return out


def _from_docx(raw: bytes) -> list[bytes]:
    from docx import Document

    document = Document(io.BytesIO(raw))
    out: list[bytes] = []
    for rel in document.part.rels.values():
        reltype = str(getattr(rel, "reltype", "") or "")
        if "image" not in reltype.lower():
            continue
        try:
            blob = _keep(bytes(rel.target_part.blob))
        except Exception:  # noqa: BLE001 — skip a broken picture part
            blob = None
        if blob is None:
            continue
        out.append(blob)
        if len(out) >= _MAX_IMAGES:
            return out
    return out
