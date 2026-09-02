from __future__ import annotations

import base64
import io
import sys
import zipfile
from pathlib import Path
from xml.sax.saxutils import escape

import pytest
from fastapi.testclient import TestClient

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "services" / "api"))
sys.path.insert(0, str(ROOT / "services" / "orchestrator"))

from app.auth import issue_test_token
from app.main import app
from app.settings import get_settings


def _xlsx_bytes(rows: list[list[str]]) -> bytes:
    shared: list[str] = []
    cells_xml: list[str] = []
    for r, row in enumerate(rows, start=1):
        parts = []
        for c, val in enumerate(row):
            idx = len(shared)
            shared.append(val)
            col = chr(ord("A") + c)
            parts.append(f'<c r="{col}{r}" t="s"><v>{idx}</v></c>')
        cells_xml.append(f'<row r="{r}">{"".join(parts)}</row>')
    sst = "".join(f"<si><t>{escape(s)}</t></si>" for s in shared)
    sheet_xml = (
        '<?xml version="1.0"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        f'<sheetData>{"".join(cells_xml)}</sheetData></worksheet>'
    )
    workbook = (
        '<?xml version="1.0"?>'
        '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        '<sheets><sheet name="课时" sheetId="1" r:id="rId1" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"/>'
        "</sheets></workbook>"
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("xl/workbook.xml", workbook)
        zf.writestr("xl/worksheets/sheet1.xml", sheet_xml)
        zf.writestr(
            "xl/sharedStrings.xml",
            '<?xml version="1.0"?>'
            '<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
            f"{sst}</sst>",
        )
    return buf.getvalue()


@pytest.fixture()
def client(tmp_path, monkeypatch):
    database = tmp_path / "edu-files.db"
    monkeypatch.setenv("PICO_DATABASE_URL", f"sqlite+aiosqlite:///{database}")
    monkeypatch.setenv("PICO_JWT_SECRET", "test-secret-at-least-32-bytes-long!!")
    monkeypatch.setenv("PICO_ENV", "development")
    from app import db as dbmod

    get_settings.cache_clear()
    dbmod._engine = None
    dbmod._Session = None
    with TestClient(app) as test_client:
        yield test_client


def _token() -> str:
    return issue_test_token(
        school_id="school-a",
        membership_id="m-edu",
        scopes=["ai:run", "ai:read"],
        settings=get_settings(),
    )


def test_post_files_requires_bearer(client: TestClient) -> None:
    res = client.post("/v1/files", json={"filename": "a.csv", "content_b64": "YQ=="})
    assert res.status_code == 401


def test_post_and_get_xlsx_excerpt(client: TestClient) -> None:
    raw = _xlsx_bytes([["班", "周课时"], ["一班", "5"]])
    token = _token()
    res = client.post(
        "/v1/files",
        headers={"authorization": f"Bearer {token}"},
        json={
            "filename": "课时.xlsx",
            "content_b64": base64.b64encode(raw).decode("ascii"),
        },
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["status"] == "ok"
    assert body["headline"] == "读到 2 行 / 2 列"
    assert body["id"]
    got = client.get(
        f"/v1/files/{body['id']}",
        headers={"authorization": f"Bearer {token}"},
    )
    assert got.status_code == 200
    assert got.json()["headline"] == "读到 2 行 / 2 列"


def test_post_pptx_persists_original_bytes(client: TestClient) -> None:
    from pptx import Presentation

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[0])
    slide.shapes.title.text = "封面"
    buf = io.BytesIO()
    deck.save(buf)
    raw = buf.getvalue()
    token = _token()
    res = client.post(
        "/v1/files",
        headers={"authorization": f"Bearer {token}"},
        json={
            "filename": "封面.pptx",
            "content_b64": base64.b64encode(raw).decode("ascii"),
        },
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["status"] == "ok"
    assert body["id"]
    assert "页" in body["headline"]


def test_post_markdown_is_workspace_readable(client: TestClient, tmp_path, monkeypatch) -> None:
    """T-AGENT-PLAIN-V1 F2: composer .md must land as utf8 Artifact the agent can read."""
    import asyncio

    from app.artifact_store import LedgerArtifactStore
    from app.auth import Principal

    token = _token()
    body_text = "年级：三年级二班。人数：42。学情：识字两极分化。\n"
    res = client.post(
        "/v1/files",
        headers={
            "authorization": f"Bearer {token}",
            "X-Conversation-Id": "convo-lesson",
        },
        json={
            "filename": "班情.md",
            "content_b64": base64.b64encode(body_text.encode("utf-8")).decode("ascii"),
        },
    )
    assert res.status_code == 200, res.text
    posted = res.json()
    assert posted["status"] == "ok"
    assert posted["id"]
    assert "三年级二班" in (posted.get("text") or "")

    store = LedgerArtifactStore(
        __import__("app.db", fromlist=["session_factory"]).session_factory(),
        conversation_id="convo-lesson",
    )
    principal = Principal(
        school_id="school-a",
        membership_id="m-edu",
        scopes=["ai:run", "ai:read"],
        iss="test",
        aud="test",
        exp=0,
        raw={},
    )

    async def _check() -> None:
        listed = await store.list(principal, limit=20)
        titles = [str(row.get("title") or "") for row in listed]
        assert "班情.md" in titles
        got = await store.read(principal, artifact_id=None, title="班情.md")
        assert got is not None
        assert got.get("content_encoding") == "utf8"
        assert "三年级二班" in str(got.get("content") or "")
        assert "42" in str(got.get("content") or "")

    fetched = client.get(
        f"/v1/files/{posted['id']}",
        headers={"authorization": f"Bearer {token}"},
    )
    assert fetched.status_code == 200, fetched.text
    assert "三年级二班" in (fetched.json().get("text") or "")

    asyncio.run(_check())


def test_foreign_membership_cannot_read(client: TestClient) -> None:
    raw = _xlsx_bytes([["a", "b"]])
    owner = _token()
    posted = client.post(
        "/v1/files",
        headers={"authorization": f"Bearer {owner}"},
        json={
            "filename": "课时.xlsx",
            "content_b64": base64.b64encode(raw).decode("ascii"),
        },
    )
    file_id = posted.json()["id"]
    other = issue_test_token(
        school_id="school-a",
        membership_id="m-other",
        scopes=["ai:read"],
        settings=get_settings(),
    )
    res = client.get(
        f"/v1/files/{file_id}",
        headers={"authorization": f"Bearer {other}"},
    )
    assert res.status_code == 404


def test_unread_pdf_still_lands_in_cabinet(client: TestClient) -> None:
    token = _token()
    headers = {"authorization": f"Bearer {token}"}
    res = client.post(
        "/v1/files",
        headers=headers,
        json={
            "filename": "坏.pdf",
            "content_b64": base64.b64encode(b"%PDF-1.4\n%not-a-real-pdf").decode("ascii"),
        },
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["id"]
    listed = client.get("/v1/artifacts", params={"mine": True}, headers=headers)
    assert listed.status_code == 200, listed.text
    rows = listed.json()["artifacts"]
    titles = [row["title"] for row in rows]
    assert "坏.pdf" in titles
    kinds = {row["kind"] for row in rows}
    assert "edu_excerpt" not in kinds
    assert "kb_text" not in kinds


def test_upload_into_folder_via_multipart(client: TestClient) -> None:
    token = _token()
    headers = {"authorization": f"Bearer {token}"}
    folder = client.post("/v1/my/folders", json={"name": "备课"}, headers=headers)
    assert folder.status_code == 200, folder.text
    folder_id = folder.json()["folder"]["id"]
    res = client.post(
        "/v1/files",
        headers=headers,
        files={"file": ("班情.md", "年级：三班\n".encode(), "text/markdown")},
        data={"folder_id": folder_id},
    )
    assert res.status_code == 200, res.text
    assert res.json()["id"]
    listed = client.get(
        "/v1/artifacts",
        params={"mine": True, "folder_id": folder_id},
        headers=headers,
    )
    assert listed.status_code == 200, listed.text
    titles = [row["title"] for row in listed.json()["artifacts"]]
    assert "班情.md" in titles
    kinds = [row["kind"] for row in listed.json()["artifacts"]]
    assert "edu_excerpt" not in kinds


def test_unknown_folder_fails_closed(client: TestClient) -> None:
    token = _token()
    res = client.post(
        "/v1/files",
        headers={"authorization": f"Bearer {token}"},
        json={
            "filename": "a.md",
            "content_b64": base64.b64encode(b"hi").decode("ascii"),
            "folder_id": "00000000-0000-4000-8000-000000000001",
        },
    )
    assert res.status_code == 404


def test_xlsx_cabinet_lists_one_row(client: TestClient) -> None:
    raw = _xlsx_bytes([["班", "周课时"], ["一班", "5"]])
    token = _token()
    headers = {"authorization": f"Bearer {token}"}
    res = client.post(
        "/v1/files",
        headers=headers,
        json={
            "filename": "课时.xlsx",
            "content_b64": base64.b64encode(raw).decode("ascii"),
        },
    )
    assert res.status_code == 200, res.text
    listed = client.get("/v1/artifacts", params={"mine": True}, headers=headers)
    rows = [row for row in listed.json()["artifacts"] if row["title"] == "课时.xlsx"]
    assert len(rows) == 1
    assert rows[0]["kind"] == "edu_office"


def test_system_tells_pi_to_read_paperclip_documents() -> None:
    system = (
        ROOT
        / "services"
        / "orchestrator"
        / "pico_orchestrator"
        / "agent_assets"
        / "system.md"
    ).read_text(encoding="utf-8")
    assert "Documents attached this turn" in system
    assert "workspace_read_file" in system
    assert "this-turn images" in system
    assert "do not say you cannot read the file" in system
    bridge = (ROOT / "services" / "true_pi_bridge" / "pico-gateway-tools.ts").read_text(
        encoding="utf-8"
    )
    assert "this-turn chat paperclip documents" in bridge


def test_inject_conversation_uploads_noop_when_empty() -> None:
    from app.edu_files import inject_conversation_uploads

    prompt = "这是什么"
    assert inject_conversation_uploads(prompt, []) == prompt
    assert inject_conversation_uploads(prompt, None) == prompt
    assert "回形针" not in inject_conversation_uploads(prompt, [])


def test_inject_conversation_uploads_cites_excerpt() -> None:
    from app.edu_files import inject_conversation_uploads

    prompt = "这是什么"
    out = inject_conversation_uploads(
        prompt,
        [
            {
                "id": "art-pdf-1",
                "title": "生物答案.pdf",
                "excerpt": "口令 PICO860-LANTERN 这是 PDF 正文",
            }
        ],
    )
    assert "生物答案.pdf" in out
    assert "PICO860-LANTERN" in out
    assert "art-pdf-1" in out
    assert out.endswith(prompt)
    assert "学校库" in out


def test_inject_scan_pdf_does_not_weld_unreadable() -> None:
    from app.edu_files import inject_conversation_uploads

    out = inject_conversation_uploads(
        "这是什么",
        [
            {
                "id": "art-scan-1",
                "title": "地理答案（XLM1）含补充说明.pdf",
                "excerpt": "",
                "error": "没抽出正文",
                "page_count": 8,
            }
        ],
    )
    assert "地理答案" in out
    assert "直接看图" in out
    assert "把相关页截图上传" not in out
    assert "另存为能复制文字" not in out
    assert "没抽出正文" not in out


def test_inject_conversation_uploads_legacy_doc_is_honest() -> None:
    from app.edu_files import inject_conversation_uploads

    out = inject_conversation_uploads(
        "看看这是什么",
        [
            {
                "id": "art-doc-1",
                "title": "教师教学计划.doc",
                "excerpt": "",
                "error": "旧版 .doc/.ppt/.xls 打不开也转不了。请另存为 .docx/.pptx/.xlsx 再试。",
            }
        ],
    )
    assert "教师教学计划.doc" in out
    assert "另存为" in out
    assert "没收到" not in out


def test_paperclip_pdf_injects_into_this_conversation(client: TestClient) -> None:
    import asyncio

    from app.auth import Principal
    from app.db import session_factory
    from app.edu_files import inject_conversation_uploads, uploads_for_conversation

    token = _token()
    body_text = "年级：三年级二班。人数：42。\n"
    res = client.post(
        "/v1/files",
        headers={
            "authorization": f"Bearer {token}",
            "X-Conversation-Id": "convo-paperclip",
        },
        json={
            "filename": "班情.md",
            "content_b64": base64.b64encode(body_text.encode("utf-8")).decode("ascii"),
        },
    )
    assert res.status_code == 200, res.text
    principal = Principal(
        school_id="school-a",
        membership_id="m-edu",
        scopes=["ai:run", "ai:read"],
        iss="test",
        aud="test",
        exp=0,
        raw={},
    )

    async def _load() -> list:
        factory = session_factory()
        async with factory() as session:
            mine = await uploads_for_conversation(session, principal, "convo-paperclip")
            other = await uploads_for_conversation(session, principal, "convo-other")
        return mine, other

    mine, other = asyncio.run(_load())
    assert other == []
    assert any("三年级二班" in str(row.get("excerpt") or "") for row in mine)
    injected = inject_conversation_uploads("这是什么", mine)
    assert "三年级二班" in injected
    assert injected.endswith("这是什么")


def test_legacy_doc_lands_unread_not_dropped(client: TestClient) -> None:
    import asyncio

    from app.auth import Principal
    from app.db import session_factory
    from app.edu_files import inject_conversation_uploads, uploads_for_conversation

    token = _token()
    res = client.post(
        "/v1/files",
        headers={
            "authorization": f"Bearer {token}",
            "X-Conversation-Id": "convo-legacy-doc",
        },
        json={
            "filename": "教师教学计划.doc",
            "content_b64": base64.b64encode(b"OLE-not-ooxml").decode("ascii"),
        },
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["id"]
    assert body["status"] != "ok"
    assert "另存为" in str(body.get("error") or body.get("headline") or "")
    principal = Principal(
        school_id="school-a",
        membership_id="m-edu",
        scopes=["ai:run", "ai:read"],
        iss="test",
        aud="test",
        exp=0,
        raw={},
    )

    async def _load() -> list:
        factory = session_factory()
        async with factory() as session:
            return await uploads_for_conversation(session, principal, "convo-legacy-doc")

    rows = asyncio.run(_load())
    assert rows
    assert rows[0]["title"] == "教师教学计划.doc"
    injected = inject_conversation_uploads("看看这是什么", rows)
    assert "另存为" in injected
    assert "教师教学计划.doc" in injected


def _visible_scan_pdf(token: str) -> bytes:
    from io import BytesIO

    from fpdf import FPDF
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (800, 1000), "white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 40)
    except OSError:
        font = ImageFont.load_default()
    draw.text((40, 80), token, fill=(0, 0, 0), font=font)
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    pdf = FPDF(unit="pt", format=(800, 1000))
    pdf.set_auto_page_break(False)
    pdf.add_page()
    pdf.image(buf, x=0, y=0, w=800, h=1000)
    return bytes(pdf.output())


def test_scan_pdf_paperclip_pages_go_to_vision(client: TestClient) -> None:
    import asyncio
    import json

    import pytest
    from sqlalchemy import select

    pytest.importorskip("pypdfium2")
    pytest.importorskip("PIL")
    pytest.importorskip("fpdf")

    from app.auth import Principal
    from app.db import ArtifactRow, session_factory
    from app.edu_files import (
        ensure_paperclip_pdf_pages,
        inject_conversation_uploads,
        uploads_for_conversation,
    )
    from pico_orchestrator.vision import clear_conversation_images, conversation_images

    token_text = "PICO860-GEO-SCAN"
    raw = _visible_scan_pdf(token_text)
    token = _token()
    cid = "convo-scan-pdf"
    clear_conversation_images(cid)
    res = client.post(
        "/v1/files",
        headers={
            "authorization": f"Bearer {token}",
            "X-Conversation-Id": cid,
        },
        json={
            "filename": "地理答案（XLM1）含补充说明.pdf",
            "content_b64": base64.b64encode(raw).decode("ascii"),
        },
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["id"]
    assert body.get("page_count") == 1
    assert "page_pngs" not in body
    pending = conversation_images(cid)
    assert pending
    assert pending[0].get("source") == "pdf-page"
    principal = Principal(
        school_id="school-a",
        membership_id="m-edu",
        scopes=["ai:run", "ai:read"],
        iss="test",
        aud="test",
        exp=0,
        raw={},
    )

    async def _round() -> tuple[list, dict]:
        factory = session_factory()
        async with factory() as session:
            rows = await uploads_for_conversation(session, principal, cid)
            excerpt = (
                await session.execute(
                    select(ArtifactRow).where(ArtifactRow.kind == "edu_excerpt")
                )
            ).scalars().first()
            sidecar = json.loads(excerpt.inline) if excerpt and excerpt.inline else {}
        return rows, sidecar

    rows, sidecar = asyncio.run(_round())
    assert "page_pngs" not in sidecar
    assert sidecar.get("page_count") == 1
    assert len(json.dumps(sidecar)) < 20_000
    injected = inject_conversation_uploads("这是什么", rows)
    assert "直接看图" in injected
    assert "把相关页截图上传" not in injected
    assert "没抽出正文" not in injected

    clear_conversation_images(cid)
    assert conversation_images(cid) == []

    async def _ensure() -> list:
        factory = session_factory()
        async with factory() as session:
            items = await uploads_for_conversation(session, principal, cid)
            await ensure_paperclip_pdf_pages(session, principal, cid, items)
            return items

    items = asyncio.run(_ensure())
    assert conversation_images(cid)
    assert conversation_images(cid)[0].get("source") == "pdf-page"
    assert int(items[0].get("page_count") or 0) >= 1
    injected_again = inject_conversation_uploads("这是什么", items)
    assert "直接看图" in injected_again
    clear_conversation_images(cid)
