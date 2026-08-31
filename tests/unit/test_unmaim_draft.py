"""T-UNMAIM-DRAFT #829: do not maim the model's HTML/PPT draft."""

from __future__ import annotations

import re
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "services" / "orchestrator"))

from pico_orchestrator.document_generators import (
    DOC_BODY_MAX,
    build_html_document,
)
from pico_orchestrator.gateway import ToolError
from pico_orchestrator.office.inspect import inspect_office_bytes
from pico_orchestrator.office.sandbox_lib import (
    STDLIB_OK,
    assert_pptx_lib_source,
    run_pptx_lib_source,
)
from pico_orchestrator.tools_builtin import build_default_gateway


def _stdlib_names(path: Path) -> set[str]:
    text = path.read_text(encoding="utf-8")
    match = re.search(r"STDLIB_OK = frozenset\(\s*\{([^}]+)\}", text)
    assert match, f"STDLIB_OK missing in {path}"
    return set(re.findall(r'"([A-Za-z0-9_]+)"', match.group(1)))


def test_html_marker_is_hidden_not_chrome() -> None:
    marker = "UNMAIM_HTML_MARK"
    full = """<!DOCTYPE html><html><head><title>t</title></head>
<body><h1>课</h1><button type="button">开始</button></body></html>"""
    text = build_html_document(title="a.html", marker=marker, body=full).decode("utf-8")
    assert marker in text
    assert 'data-pico-marker="' + marker + '"' in text
    assert "标记：" not in text
    assert "data-pico-marker-line" not in text
    prose = build_html_document(title="p.html", marker=marker, body="一段说明。").decode(
        "utf-8"
    )
    assert marker in prose
    assert "标记：" not in prose


def test_html_over_50k_lands_with_tail_script() -> None:
    sentinel = "PICO_UNMAIM_TAIL_OK"
    pad = "x" * 55_000
    body = f"""<!DOCTYPE html>
<html lang="zh-CN"><head><meta charset="utf-8"/><title>long</title></head>
<body>
<canvas id="c" width="80" height="60"></canvas>
<!-- {pad} -->
<script>window.__picoUnmaim = "{sentinel}";</script>
</body></html>"""
    assert len(body) > 50_000
    assert len(body) < DOC_BODY_MAX
    text = build_html_document(title="long.html", marker="LONG1", body=body).decode(
        "utf-8"
    )
    assert sentinel in text
    assert "<canvas" in text.lower()
    assert "标记：" not in text


def test_html_over_max_fails_closed_no_slice() -> None:
    body = "<!DOCTYPE html><html><body>" + ("y" * (DOC_BODY_MAX + 1)) + "</body></html>"
    with pytest.raises(ValueError, match="不会截断"):
        build_html_document(title="over.html", marker="OVER1", body=body)


def test_html_cdn_still_fail_closed() -> None:
    body = """<!DOCTYPE html><html><body>
<script type="module">import * as THREE from "https://cdn.jsdelivr.net/npm/three";</script>
</body></html>"""
    with pytest.raises(ValueError, match="外网资源"):
        build_html_document(title="cdn.html", marker="CDN1", body=body)


def test_system_does_not_steer_canvas_downgrade() -> None:
    system = (
        ROOT
        / "services"
        / "orchestrator"
        / "pico_orchestrator"
        / "agent_assets"
        / "system.md"
    ).read_text(encoding="utf-8")
    assert "rewrite the page with canvas drawing" not in system
    assert "dumb the page down" in system
    assert "copy" in system and "BytesIO" in system


def test_gateway_html_tool_does_not_require_canvas_only() -> None:
    ts = (
        ROOT / "services" / "true_pi_bridge" / "pico-gateway-tools.ts"
    ).read_text(encoding="utf-8")
    assert "canvas only" not in ts
    assert "canvas allowed, not required" in ts
    assert "do not dumb it down" in ts


def test_pptx_stdlib_allowlist_matches_exec() -> None:
    lib = ROOT / "services/orchestrator/pico_orchestrator/office/sandbox_lib.py"
    exe = ROOT / "services/orchestrator/pico_orchestrator/office/sandbox_exec.py"
    names = _stdlib_names(lib)
    assert names == _stdlib_names(exe)
    assert names == set(STDLIB_OK)
    assert "copy" in STDLIB_OK
    assert "os" not in STDLIB_OK
    assert "sys" not in STDLIB_OK


def test_pptx_naked_gpt_stdlib_lands() -> None:
    source = """
from pptx import Presentation, Inches
from pptx.dml.color import RGBColor
import copy
import math
from datetime import datetime
from io import BytesIO

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(9), Inches(1.2))
box.text_frame.text = f"ok {math.pi:.2f} {datetime(2026, 8, 31):%Y-%m-%d} {copy.copy([1])[0]}"
buf = BytesIO()
assert buf.getvalue() == b""
save_deck(prs)
"""
    assert_pptx_lib_source(source)
    raw = run_pptx_lib_source(source)
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) >= 1
    units = outline.get("units") or []
    blob = " ".join(
        " ".join(str(b) for b in (u.get("bullets") or [])) for u in units if isinstance(u, dict)
    )
    assert "ok" in blob
    assert "2026-08-31" in blob


def test_pptx_io_open_denied() -> None:
    with pytest.raises(ToolError) as ei:
        assert_pptx_lib_source(
            "from io import open\nfrom pptx import Presentation\nsave_deck(Presentation())"
        )
    assert ei.value.code == "sandbox.exec_denied"


def test_generate_html_tool_accepts_over_50k() -> None:
    import asyncio
    from typing import Any, ClassVar

    from pico_orchestrator.gateway import Principal

    class _P:
        school_id = "s1"
        membership_id = "m1"
        scopes: ClassVar[list[str]] = ["ai:run"]

    class _MemStore:
        def __init__(self) -> None:
            self.items: dict[str, dict[str, Any]] = {}

        async def write(
            self, principal: Principal, *, title: str, content: str | bytes, kind: str
        ) -> dict[str, Any]:
            aid = f"a-{len(self.items)+1}"
            body = content if isinstance(content, str) else content.decode("utf-8", "replace")
            self.items[aid] = {
                "artifact_id": aid,
                "title": title,
                "kind": kind,
                "content": body,
            }
            return dict(self.items[aid])

        async def read(
            self,
            principal: Principal,
            *,
            artifact_id: str | None,
            title: str | None,
        ) -> dict[str, Any] | None:
            if artifact_id and artifact_id in self.items:
                return self.items[artifact_id]
            return None

        async def list(self, principal: Principal, *, limit: int) -> list[dict[str, Any]]:
            return list(self.items.values())[:limit]

    sentinel = "PICO_TOOL_TAIL_OK"
    pad = "z" * 55_000
    body = (
        "<!DOCTYPE html><html><head><title>t</title></head><body>"
        f"<p>ok</p><!-- {pad} --><script>window.__t='{sentinel}';</script>"
        "</body></html>"
    )
    assert len(body) > 50_000
    store = _MemStore()
    gw = build_default_gateway(store)

    async def _run() -> dict[str, Any]:
        return await gw.invoke(
            _P(),  # type: ignore[arg-type]
            "generate_html_document",
            {"title": "long.html", "marker": "TOOL1", "body": body},
        )

    out = asyncio.run(_run())
    aid = str(out.get("artifact_id") or "")
    assert aid
    landed = store.items[aid]["content"]
    assert sentinel in landed
    assert "标记：" not in landed
