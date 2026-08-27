"""Default Title-and-Content well: body/picture must not collapse over the title."""

from __future__ import annotations

import io

from pico_orchestrator.office.render import render_spec
from pico_orchestrator.office.spec import parse_spec
from pico_orchestrator.sandbox_s2 import encode_rgb_png
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

WIDE_PNG = encode_rgb_png(320, 180, bytes((30, 120, 200)) * (320 * 180))
TALL_PNG = encode_rgb_png(120, 800, bytes((200, 80, 30)) * (120 * 800))
SLIDE_H = int(Inches(7.5))
TITLE_BOTTOM_MIN = int(Inches(1.4))


def _slide0(raw: bytes):
    return Presentation(io.BytesIO(raw)).slides[0]


def _placeholder_idx(shape) -> int | None:
    try:
        return int(shape.placeholder_format.idx)
    except (AttributeError, ValueError):
        return None


def _parts(slide) -> tuple[object, object | None, object | None]:
    title = slide.shapes.title
    body = None
    pic = None
    for shape in slide.shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            pic = shape
            continue
        idx = _placeholder_idx(shape)
        if idx == 0:
            title = shape
        elif idx == 1:
            body = shape
    return title, body, pic


def test_image_and_bullets_body_has_real_height_below_title() -> None:
    spec = parse_spec(
        {
            "kind": "pptx",
            "title": "井",
            "blocks": [
                {
                    "type": "slide",
                    "title": "为什么选择豌豆？",
                    "bullets": ["自花传粉", "性状明显", "周期短"],
                    "image_artifact_id": "img-1",
                }
            ],
        }
    )
    raw = render_spec(spec, images={"img-1": WIDE_PNG})
    title, body, pic = _parts(_slide0(raw))
    assert body is not None and pic is not None
    assert int(body.height) >= int(Inches(3.0))
    assert int(body.top) >= TITLE_BOTTOM_MIN
    assert int(title.top) + int(title.height) <= int(body.top) + int(Inches(0.15))
    assert int(pic.top) >= TITLE_BOTTOM_MIN
    assert int(pic.top) + int(pic.height) <= SLIDE_H
    assert int(pic.left) >= int(body.left) + int(body.width) - int(Inches(0.2))


def test_tall_picture_is_capped_inside_slide() -> None:
    spec = parse_spec(
        {
            "kind": "pptx",
            "title": "高图",
            "blocks": [
                {
                    "type": "slide",
                    "title": "流程",
                    "bullets": ["P", "F1", "F2"],
                    "image_artifact_id": "tall",
                }
            ],
        }
    )
    raw = render_spec(spec, images={"tall": TALL_PNG})
    _title, _body, pic = _parts(_slide0(raw))
    assert pic is not None
    assert int(pic.height) <= int(Inches(5.3))
    assert int(pic.top) + int(pic.height) <= SLIDE_H


def test_text_only_slide_keeps_default_body_box() -> None:
    spec = parse_spec(
        {
            "kind": "pptx",
            "title": "目录",
            "blocks": [{"type": "slide", "title": "目录", "bullets": ["一", "二", "三"]}],
        }
    )
    raw = render_spec(spec, images={})
    title, body, pic = _parts(_slide0(raw))
    assert pic is None and body is not None
    assert int(body.top) >= TITLE_BOTTOM_MIN
    assert int(body.height) >= int(Inches(4.0))
    assert int(title.top) + int(title.height) <= int(body.top) + int(Inches(0.15))
