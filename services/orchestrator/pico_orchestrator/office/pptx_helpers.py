"""Thin python-pptx helpers for the isolated sandbox.

Not a second Office OS. These are the same verbs a document-skill
script already writes by hand: title slide, bullets, table.
Upstream: python-pptx. No new spec fields.
"""

from __future__ import annotations

from typing import Any


class ImagePathMap(dict):
    """Ledger images: dict[id]=path, and IMAGE_PATHS[0] is the first path."""

    def __getitem__(self, key: Any) -> Any:
        if isinstance(key, int) and not isinstance(key, bool):
            return list(self.values())[key]
        if key in self:
            return super().__getitem__(key)
        text = str(key)
        if text in self:
            return super().__getitem__(text)
        if isinstance(key, str) and key.isdigit():
            return list(self.values())[int(key)]
        return super().__getitem__(key)

    def __contains__(self, key: object) -> bool:
        if isinstance(key, int) and not isinstance(key, bool):
            n = len(self)
            return -n <= key < n if n else False
        return super().__contains__(key)

    def get(self, key: Any, default: Any = None) -> Any:  # type: ignore[override]
        try:
            return self[key]
        except (KeyError, IndexError, TypeError):
            return default


def add_title_slide(
    prs: Any,
    title: str = "",
    subtitle: str = "",
    image: Any = None,
    picture: Any = None,
    image_path: Any = None,
) -> Any:
    """Title layout (owner/date belong in subtitle). ``image=`` embeds a picture."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title_shape.text = str(title or "")
    if subtitle:
        _set_placeholder_text(slide, 1, str(subtitle))
    pic = _first_image(image, picture, image_path)
    if pic is not None:
        _embed_picture(slide, pic, cover=True)
    return slide


def add_content_slide(
    prs: Any,
    title: str = "",
    bullets: Any = (),
    image: Any = None,
    picture: Any = None,
    image_path: Any = None,
) -> Any:
    """Title-and-content layout. ``bullets`` is a list of strings."""
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title_shape.text = str(title or "")
    items = [str(item).strip() for item in (bullets or ()) if str(item).strip()]
    if items:
        _set_placeholder_text(slide, 1, "\n".join(items))
    pic = _first_image(image, picture, image_path)
    if pic is not None:
        _embed_picture(slide, pic, cover=False)
    return slide


def add_table(
    slide: Any = None,
    rows: Any = None,
    left: Any = None,
    top: Any = None,
    width: Any = None,
    height: Any = None,
    *,
    prs: Any = None,
    presentation: Any = None,
    data: Any = None,
    cells: Any = None,
    headers: Any = None,
    title: str = "",
) -> Any:
    """Add a table from a list of row lists. Sizes default to the content well.

    Aliases: ``prs=`` / ``presentation=`` (add a slide then the table),
    ``data=`` / ``cells=`` for the grid, ``headers=`` prepended as row 0.
    A Presentation passed as the first positional is treated as ``prs``.
    """
    from pptx.util import Inches, Pt

    deck = prs if prs is not None else presentation
    if deck is None and _is_presentation(slide):
        deck = slide
        slide = None
    grid_src = rows if rows is not None else data if data is not None else cells
    if headers is not None:
        header_row = ["" if cell is None else str(cell) for cell in headers]
        body = list(grid_src or [])
        grid_src = [header_row, *body]
    if grid_src is None:
        raise TypeError("add_table 需要 rows（或 data/cells）")
    if slide is None:
        if deck is None:
            raise TypeError("add_table 需要 slide 或 prs")
        slide = add_content_slide(deck, str(title or ""), ())
    grid = _table_grid(grid_src)
    n_rows = len(grid)
    n_cols = len(grid[0])
    shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        left if left is not None else Inches(0.5),
        top if top is not None else Inches(1.7),
        width if width is not None else Inches(9.0),
        height if height is not None else Inches(min(5.2, 0.42 * n_rows + 0.5)),
    )
    table = shape.table
    for r_i, row in enumerate(grid):
        for c_i, cell in enumerate(row):
            table.cell(r_i, c_i).text = cell
            if r_i == 0:
                for para in table.cell(r_i, c_i).text_frame.paragraphs:
                    for run in para.runs:
                        run.font.bold = True
                        run.font.size = Pt(14)
    return shape


def pipe_table_rows(bullets: Any) -> tuple[tuple[str, ...], ...] | None:
    """If every bullet is ``a|b`` (or more), return a table grid. Else None."""
    items = [str(item).strip() for item in (bullets or ()) if str(item).strip()]
    if len(items) < 2:
        return None
    rows: list[tuple[str, ...]] = []
    for item in items:
        if "|" not in item:
            return None
        cells = [part.strip() for part in item.split("|")]
        if cells and cells[0] == "":
            cells = cells[1:]
        if cells and cells[-1] == "":
            cells = cells[:-1]
        if len(cells) < 2:
            return None
        rows.append(tuple(cells))
    width = max(len(row) for row in rows)
    if width < 2:
        return None
    return tuple(row + ("",) * (width - len(row)) for row in rows)


def _table_grid(rows: Any) -> list[list[str]]:
    if not rows:
        raise ValueError("add_table 需要至少一行。")
    grid: list[list[str]] = []
    width = 0
    for row in rows:
        if isinstance(row, (str, bytes)):
            raise TypeError("add_table 每一行必须是单元格列表。")
        cells = ["" if cell is None else str(cell) for cell in row]
        if not cells:
            raise ValueError("add_table 每一行必须是非空列表。")
        width = max(width, len(cells))
        grid.append(cells)
    if width < 1:
        raise ValueError("add_table 需要至少一列。")
    return [row + [""] * (width - len(row)) for row in grid]


def _set_placeholder_text(slide: Any, idx: int, text: str) -> None:
    placeholders = getattr(slide, "placeholders", None)
    if placeholders is None:
        return
    for shape in placeholders:
        fmt = getattr(shape, "placeholder_format", None)
        if getattr(fmt, "idx", None) == idx and getattr(shape, "has_text_frame", False):
            shape.text = text
            return


def _is_presentation(obj: Any) -> bool:
    return (
        obj is not None
        and hasattr(obj, "slides")
        and hasattr(obj, "slide_layouts")
        and not hasattr(obj, "shapes")
    )


def _first_image(*candidates: Any) -> Any:
    for item in candidates:
        if item is None or item is False:
            continue
        if isinstance(item, str) and not item.strip():
            continue
        return item
    return None


def _embed_picture(slide: Any, image: Any, *, cover: bool) -> None:
    """Thin add_picture. Missing file skips (same as missing image_artifact_id)."""
    from io import BytesIO
    from pathlib import Path

    from pptx.util import Inches

    if isinstance(image, (bytes, bytearray)):
        target: Any = BytesIO(image)
    elif isinstance(image, (str, Path)):
        path = Path(image)
        if not path.is_file():
            return
        target = str(path)
    else:
        target = image
    left = Inches(1.2 if cover else 5.15)
    top = Inches(3.2 if cover else 1.7)
    width = Inches(7.6 if cover else 4.5)
    try:
        slide.shapes.add_picture(target, left, top, width=width)
    except OSError:
        return
