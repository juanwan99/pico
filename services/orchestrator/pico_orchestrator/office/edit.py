"""Path B: inspect address → thin in-place edit. Upstream python-docx / python-pptx.

Existing ``edit_docx_bytes`` / ``edit_pptx_title_bytes`` stay as aliases.
"""

from __future__ import annotations

import io
import re

from pico_orchestrator.artifact_types import is_valid_ooxml_package

_P = re.compile(r"^p:(\d+)$", re.IGNORECASE)
_CELL = re.compile(r"^t:(\d+)\.r(\d+)\.c(\d+)$", re.IGNORECASE)
_SLIDE_TITLE = re.compile(r"^s:(\d+)(?:\.title)?$", re.IGNORECASE)
_SLIDE_BULLET = re.compile(r"^s:(\d+)\.b:(\d+)$", re.IGNORECASE)


def edit_docx_bytes(
    raw: bytes,
    *,
    paragraph_index: int,
    text: str,
) -> bytes:
    """Replace one 1-based nonempty paragraph. Other paragraphs stay."""
    return edit_by_address(raw, ext=".docx", address=f"p:{paragraph_index}", text=text)


def edit_pptx_title_bytes(
    raw: bytes,
    *,
    slide_index: int,
    new_title: str,
) -> bytes:
    """Replace the title of one 1-based slide. Other slides stay."""
    return edit_by_address(raw, ext=".pptx", address=f"s:{slide_index}.title", text=new_title)


def edit_by_address(
    raw: bytes,
    *,
    ext: str,
    address: str,
    text: str,
) -> bytes:
    """Path B: change one address. Never rebuild the package from a guessed spec."""
    suffix = (ext or "").lower()
    if suffix and not suffix.startswith("."):
        suffix = f".{suffix}"
    new_text = (text or "").strip()
    if not new_text:
        raise ValueError("新内容不能为空。")
    addr = (address or "").strip()
    if not addr:
        raise ValueError("请给出地址（inspect_document 返回的 addr）。")
    if suffix == ".docx":
        return _edit_docx(raw, addr, new_text)
    if suffix == ".pptx":
        return _edit_pptx(raw, addr, new_text)
    raise ValueError("只支持改 .docx / .pptx。")


def _edit_docx(raw: bytes, address: str, new_text: str) -> bytes:
    if not is_valid_ooxml_package(raw, ".docx"):
        raise ValueError("不是真 Word（OOXML）。请上传已有的 .docx。")
    from docx import Document

    document = Document(io.BytesIO(raw))
    m_p = _P.match(address)
    m_cell = _CELL.match(address)
    if m_p:
        index = int(m_p.group(1))
        paras = [p for p in document.paragraphs if str(getattr(p, "text", "") or "").strip()]
        if not paras:
            raise ValueError("这份 Word 没有可改的段落。")
        if index < 1 or index > len(paras):
            raise ValueError(f"没有第 {index} 段（共 {len(paras)} 段）。")
        paras[index - 1].text = new_text
    elif m_cell:
        t_i, r_i, c_i = (int(m_cell.group(1)), int(m_cell.group(2)), int(m_cell.group(3)))
        tables = list(document.tables)
        if t_i < 1 or t_i > len(tables):
            raise ValueError(f"没有第 {t_i} 张表（共 {len(tables)} 张）。")
        table = tables[t_i - 1]
        rows = list(table.rows)
        if r_i < 1 or r_i > len(rows):
            raise ValueError(f"表 {t_i} 没有第 {r_i} 行。")
        cells = list(rows[r_i - 1].cells)
        if c_i < 1 or c_i > len(cells):
            raise ValueError(f"表 {t_i} 第 {r_i} 行没有第 {c_i} 列。")
        cells[c_i - 1].text = new_text
    else:
        raise ValueError(f"不认识的 Word 地址：{address}。用 p:N 或 t:N.rR.cC。")
    out = io.BytesIO()
    document.save(out)
    data = out.getvalue()
    if not is_valid_ooxml_package(data, ".docx"):
        raise ValueError("改完后不是真 Word，未保存。")
    return data


def _edit_pptx(raw: bytes, address: str, new_text: str) -> bytes:
    if not is_valid_ooxml_package(raw, ".pptx"):
        raise ValueError("不是真 PPT（OOXML）。请上传已有的 .pptx。")
    from pptx import Presentation

    deck = Presentation(io.BytesIO(raw))
    slides = list(deck.slides)
    if not slides:
        raise ValueError("这份 PPT 没有可改的页。")
    m_title = _SLIDE_TITLE.match(address)
    m_bullet = _SLIDE_BULLET.match(address)
    if m_bullet:
        s_i, b_i = int(m_bullet.group(1)), int(m_bullet.group(2))
        if s_i < 1 or s_i > len(slides):
            raise ValueError(f"没有第 {s_i} 页（共 {len(slides)} 页）。")
        _set_slide_bullet(slides[s_i - 1], b_i, new_text)
    elif m_title:
        s_i = int(m_title.group(1))
        if s_i < 1 or s_i > len(slides):
            raise ValueError(f"没有第 {s_i} 页（共 {len(slides)} 页）。")
        _set_slide_title(slides[s_i - 1], new_text)
    else:
        raise ValueError(f"不认识的 PPT 地址：{address}。用 s:N.title 或 s:N.b:M。")
    out = io.BytesIO()
    deck.save(out)
    data = out.getvalue()
    if not is_valid_ooxml_package(data, ".pptx"):
        raise ValueError("改完后不是真 PPT，未保存。")
    return data


def _title_shape(slide: object) -> object | None:
    shape = getattr(slide.shapes, "title", None)
    if shape is not None:
        return shape
    for item in slide.shapes:
        idx = getattr(getattr(item, "placeholder_format", None), "idx", None)
        if idx == 0 and getattr(item, "has_text_frame", False):
            return item
    return None


def _set_slide_title(slide: object, title: str) -> None:
    shape = _title_shape(slide)
    if shape is not None and getattr(shape, "has_text_frame", False):
        shape.text = title
        return
    for item in slide.shapes:
        if getattr(item, "has_text_frame", False):
            item.text = title
            return
    raise ValueError("该页没有标题可改。")


def _set_slide_bullet(slide: object, bullet_index: int, text: str) -> None:
    title_shape = _title_shape(slide)
    for shape in slide.shapes:
        idx = getattr(getattr(shape, "placeholder_format", None), "idx", None)
        name = str(getattr(shape, "name", "") or "").lower()
        if shape is title_shape or idx == 0 or name.startswith("title"):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        paras = [p for p in shape.text_frame.paragraphs if str(p.text or "").strip()]
        if bullet_index < 1 or bullet_index > len(paras):
            raise ValueError(f"该页没有第 {bullet_index} 条要点（共 {len(paras)} 条）。")
        paras[bullet_index - 1].text = text
        return
    raise ValueError("该页没有要点可改。")
