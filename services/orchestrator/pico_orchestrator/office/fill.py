"""Replace {{key}} placeholders. Missing keys stay; leftover is reported."""

from __future__ import annotations

import io
import re
from dataclasses import dataclass

from pico_orchestrator.artifact_types import is_valid_ooxml_package
from pico_orchestrator.office.legacy import require_supported_office_ext

PLACEHOLDER_RE = re.compile(r"\{\{\s*([^{}]+?)\s*\}\}")


def fill_placeholders(
    text: str, values: dict[str, str], hits: list[str] | None = None
) -> str:
    if not text or "{{" not in text:
        return text

    def _one(match: re.Match[str]) -> str:
        key = match.group(1).strip()
        if key in values:
            if hits is not None:
                hits.append(key)
            return str(values[key])
        return match.group(0)

    return PLACEHOLDER_RE.sub(_one, text)


def leftover_placeholders(text: str) -> list[str]:
    if not text or "{{" not in text:
        return []
    return [m.group(1).strip() for m in PLACEHOLDER_RE.finditer(text) if m.group(1).strip()]


@dataclass(frozen=True)
class FillReceipt:
    """Honest fill result. filled is true only when at least one {{key}} hit."""

    data: bytes
    filled_keys: tuple[str, ...]
    leftover: tuple[str, ...]

    @property
    def filled(self) -> bool:
        return bool(self.filled_keys)


def fill_office_with_receipt(
    raw: bytes, ext: str, values: dict[str, str]
) -> FillReceipt:
    """Fill then report which keys actually replaced a leftover {{key}}."""
    from pico_orchestrator.office.inspect import inspect_office_bytes

    suffix = require_supported_office_ext(ext)
    hits: list[str] = []
    data = fill_office_bytes(raw, suffix, values, hits=hits)
    leftover = tuple(
        inspect_office_bytes(data, suffix).get("placeholders") or []
    )
    hit_set = set(hits)
    filled_keys = tuple(key for key in values if key in hit_set)
    return FillReceipt(data=data, filled_keys=filled_keys, leftover=leftover)


def fill_office_bytes(
    raw: bytes,
    ext: str,
    values: dict[str, str],
    hits: list[str] | None = None,
) -> bytes:
    suffix = require_supported_office_ext(ext)
    if not values:
        raise ValueError("values 不能为空。请给出要替换的 {{key}}。")
    if suffix == ".docx":
        return _fill_docx(raw, values, hits)
    if suffix == ".pptx":
        return _fill_pptx(raw, values, hits)
    return _fill_xlsx(raw, values, hits)


def _fill_docx(
    raw: bytes, values: dict[str, str], hits: list[str] | None = None
) -> bytes:
    if not is_valid_ooxml_package(raw, ".docx"):
        raise ValueError("不是真 Word（OOXML）。请上传已有的 .docx。")
    from docx import Document

    document = Document(io.BytesIO(raw))
    for para in document.paragraphs:
        text = str(para.text or "")
        filled = fill_placeholders(text, values, hits)
        if filled != text:
            para.text = filled
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                text = str(cell.text or "")
                filled = fill_placeholders(text, values, hits)
                if filled != text:
                    cell.text = filled
    out = io.BytesIO()
    document.save(out)
    data = out.getvalue()
    if not is_valid_ooxml_package(data, ".docx"):
        raise ValueError("套模板后不是真 Word，未保存。")
    return data


def _fill_pptx(
    raw: bytes, values: dict[str, str], hits: list[str] | None = None
) -> bytes:
    if not is_valid_ooxml_package(raw, ".pptx"):
        raise ValueError("不是真 PPT（OOXML）。请上传已有的 .pptx。")
    from pptx import Presentation

    deck = Presentation(io.BytesIO(raw))
    for slide in deck.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = str(getattr(shape, "text", "") or "")
            filled = fill_placeholders(text, values, hits)
            if filled != text:
                shape.text = filled
    out = io.BytesIO()
    deck.save(out)
    data = out.getvalue()
    if not is_valid_ooxml_package(data, ".pptx"):
        raise ValueError("套模板后不是真 PPT，未保存。")
    return data


def _fill_xlsx(
    raw: bytes, values: dict[str, str], hits: list[str] | None = None
) -> bytes:
    if not is_valid_ooxml_package(raw, ".xlsx"):
        raise ValueError("不是真 Excel（OOXML）。请上传已有的 .xlsx。")
    from openpyxl import load_workbook

    book = load_workbook(io.BytesIO(raw))
    for sheet in book.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    filled = fill_placeholders(cell.value, values, hits)
                    if filled != cell.value:
                        cell.value = filled
    out = io.BytesIO()
    book.save(out)
    data = out.getvalue()
    if not is_valid_ooxml_package(data, ".xlsx"):
        raise ValueError("套模板后不是真 Excel，未保存。")
    return data
