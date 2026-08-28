"""S2 T-OFFICE-THICK: missing image does not fail the deck; cover/table; sandbox helpers."""

from __future__ import annotations

import io
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "services" / "orchestrator"))

from pico_orchestrator.office.inspect import inspect_office_bytes
from pico_orchestrator.office.pptx_helpers import pipe_table_rows
from pico_orchestrator.office.render import render_spec
from pico_orchestrator.office.sandbox_lib import run_pptx_lib_source
from pico_orchestrator.office.spec import parse_spec
from pico_orchestrator.sandbox_s2 import encode_rgb_png
from pico_orchestrator.tools_builtin import build_default_gateway

PNG = encode_rgb_png(64, 48, bytes((20, 80, 160)) * (64 * 48))


@dataclass
class P:
    school_id: str = "school-a"
    membership_id: str = "member-a"
    scopes: list[str] | None = None

    def __post_init__(self) -> None:
        if self.scopes is None:
            self.scopes = ["ai:run"]


class MemoryArtifactStore:
    def __init__(self) -> None:
        self.rows: list[dict[str, Any]] = []

    async def write(
        self,
        principal: P,
        *,
        title: str,
        content: str | bytes,
        kind: str,
    ) -> dict[str, Any]:
        aid = f"art-{len(self.rows) + 1}"
        if isinstance(content, bytes):
            row = {
                "artifact_id": aid,
                "title": title,
                "content": content,
                "kind": kind,
                "byte_size": len(content),
            }
        else:
            row = {
                "artifact_id": aid,
                "title": title,
                "content": content,
                "kind": kind,
                "byte_size": len(content.encode("utf-8")),
            }
        self.rows.append(row)
        return {k: v for k, v in row.items() if k != "content"}

    async def read(
        self,
        principal: P,
        *,
        artifact_id: str | None,
        title: str | None,
    ) -> dict[str, Any] | None:
        for row in reversed(self.rows):
            if artifact_id and row["artifact_id"] == artifact_id:
                return row
            if not artifact_id and title and row["title"] == title:
                return row
        return None

    async def list(self, principal: P, *, limit: int) -> list[dict[str, Any]]:
        return list(reversed(self.rows))[:limit]


def test_pipe_table_rows_needs_two_columns() -> None:
    assert pipe_table_rows(["收入|88", "同比|-12%"]) == (("收入", "88"), ("同比", "-12%"))
    assert pipe_table_rows(["只有一列"]) is None
    assert pipe_table_rows(["a|b"]) is None


def test_empty_blocks_with_title_become_one_slide() -> None:
    """Live F4: spec.blocks empty must not fail the whole PPT when title exists."""
    spec = parse_spec({"kind": "pptx", "title": "办公尺752.pptx", "blocks": []})
    assert len(spec.blocks) == 1
    assert spec.blocks[0].type == "slide"
    assert "办公尺752" in spec.blocks[0].title
    raw = render_spec(spec)
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) == 1


def test_empty_blocks_without_title_or_theme_still_fail() -> None:
    with pytest.raises(ValueError, match="不能为空"):
        parse_spec({"kind": "pptx"})


@pytest.mark.asyncio
async def test_generate_pptx_sibling_blocks_not_dropped_by_stub_spec() -> None:
    """Live F4 r1/r2: spec={images:[]} + top-level blocks was 不能为空."""
    gw = build_default_gateway(MemoryArtifactStore())
    out = await gw.invoke(
        P(),
        "generate_pptx_document",
        {
            "title": "办公尺752-f3m-r2.pptx",
            "spec": {"images": []},
            "blocks": [
                {
                    "type": "cover",
                    "title": "办公简报：本周经营风险与下周动作",
                    "bullets": ["责任人：张三", "日期：2026-08-28"],
                },
                {
                    "type": "content",
                    "title": "本周经营风险总览",
                    "bullets": ["市场端询盘下降", "交付延迟", "回款压力"],
                },
                {
                    "type": "content",
                    "title": "下周动作",
                    "bullets": ["催收回款", "锁定料源", "调整排期"],
                },
            ],
        },
    )
    assert out.get("ok") is not False
    assert out.get("format") == "pptx"
    assert int(out.get("observation", {}).get("outline", {}).get("slides") or 0) == 3


@pytest.mark.asyncio
async def test_generate_pptx_kpi_stub_spec_keeps_sibling_blocks() -> None:
    """Live F4 r1: spec={kpi_table_title} + sibling blocks."""
    gw = build_default_gateway(MemoryArtifactStore())
    out = await gw.invoke(
        P(),
        "generate_pptx_document",
        {
            "title": "办公尺752-f3m-r1.pptx",
            "spec": {"kpi_table_title": "核心数字"},
            "blocks": [
                {
                    "type": "cover",
                    "title": "本周经营风险与下周动作",
                    "bullets": ["责任人：张三"],
                },
                {
                    "type": "content",
                    "title": "风险总览",
                    "bullets": ["收入端延期", "毛利承压"],
                },
            ],
        },
    )
    assert out.get("format") == "pptx"
    assert int(out.get("observation", {}).get("outline", {}).get("slides") or 0) == 2


def test_cover_content_aliases_are_slides() -> None:
    """Live F: Pi sent type=cover/content and the whole deck died."""
    spec = parse_spec(
        {
            "kind": "pptx",
            "title": "办公尺752.pptx",
            "blocks": [
                {
                    "type": "cover",
                    "title": "本周经营风险与下周动作",
                    "bullets": ["责任人：张三", "日期：2026-08-28"],
                    "image_artifact_id": "cover-id",
                },
                {
                    "type": "content",
                    "title": "本周经营风险总览",
                    "bullets": ["收入端延期", "毛利承压", "回款变慢"],
                },
                {
                    "type": "title",
                    "title": "下周动作",
                    "bullets": ["锁定供应商", "催收回款"],
                },
            ],
        }
    )
    assert all(block.type == "slide" for block in spec.blocks)
    raw = render_spec(spec, images={})
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) == 3


@pytest.mark.asyncio
async def test_generate_pptx_live_cover_content_shape_writes() -> None:
    """Exact first-fail argument shape from run b8763b35 seq 11."""
    gw = build_default_gateway(MemoryArtifactStore())
    out = await gw.invoke(
        P(),
        "generate_pptx_document",
        {
            "title": "风险二：毛利率承压",
            "blocks": [
                {
                    "type": "cover",
                    "title": "本周经营风险与下周动作",
                    "bullets": ["责任人：张三", "日期：2026-08-28"],
                    "image_artifact_id": "a4f6458a-dcd3-4b40-bebf-9d3df2864b37",
                },
                {
                    "type": "content",
                    "title": "本周经营风险总览",
                    "bullets": [
                        "收入端：核心客户订单交付延期，风险敞口约 12%",
                        "毛利端：原材料涨价挤压毛利率约 2.3 个点",
                    ],
                },
                {
                    "type": "content",
                    "title": "风险一：订单交付延期",
                    "bullets": ["大客户 A 项目关键原料缺货，预计延期 5–7 天"],
                },
            ],
        },
    )
    assert out.get("format") == "pptx"
    assert int(out.get("observation", {}).get("outline", {}).get("slides") or 0) == 3
    assert out.get("ok") is not False


def test_missing_image_id_still_writes_deck() -> None:
    spec = parse_spec(
        {
            "kind": "pptx",
            "title": "缺图仍交.pptx",
            "blocks": [
                {
                    "type": "slide",
                    "title": "封面",
                    "bullets": ["管理层 · 2026-08-28"],
                    "image_artifact_id": "missing-cover",
                },
                {
                    "type": "slide",
                    "title": "要点",
                    "bullets": ["汇率敞口", "客户流失", "库存积压"],
                },
            ],
        }
    )
    raw = render_spec(spec, images={})
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) == 2
    assert int(outline["images"]) == 0


@pytest.mark.asyncio
async def test_generate_pptx_skips_missing_image_id() -> None:
    gw = build_default_gateway(MemoryArtifactStore())
    out = await gw.invoke(
        P(),
        "generate_pptx_document",
        {
            "title": "首写不毁.pptx",
            "blocks": [
                {
                    "type": "slide",
                    "title": "决策会",
                    "bullets": ["责任人 张三 · 2026-08-28"],
                    "image_artifact_id": "not-yet",
                },
                {"type": "slide", "title": "风险", "bullets": ["汇率", "客户", "库存"]},
            ],
        },
    )
    assert out.get("format") == "pptx"
    assert out.get("observation", {}).get("outline", {}).get("slides") == 2
    assert out.get("observation", {}).get("outline", {}).get("images") == 0
    assert "image_artifact_id" in (out.get("observation", {}).get("outline", {}).get("hint") or "")


def test_cover_uses_title_layout_and_embeds() -> None:
    spec = parse_spec(
        {
            "kind": "pptx",
            "title": "封面.pptx",
            "blocks": [
                {
                    "type": "slide",
                    "title": "管理层决策会",
                    "bullets": ["责任人 张三 · 2026-08-28"],
                    "image_artifact_id": "cover",
                }
            ],
        }
    )
    raw = render_spec(spec, images={"cover": PNG})
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["images"]) >= 1
    slide = Presentation(io.BytesIO(raw)).slides[0]
    pics = [s for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE]
    assert pics
    blob = "\n".join(
        (u.get("title") or "") + "\n" + "\n".join(u.get("bullets") or [])
        for u in (outline.get("units") or [])
        if isinstance(u, dict)
    )
    assert "张三" in blob or "2026-08-28" in blob


def test_pipe_bullets_render_as_table() -> None:
    spec = parse_spec(
        {
            "kind": "pptx",
            "title": "表.pptx",
            "blocks": [
                {
                    "type": "slide",
                    "title": "数字",
                    "bullets": ["项|本季|同比", "收入|88|-12%", "毛利|21|-4pt"],
                }
            ],
        }
    )
    raw = render_spec(spec)
    deck = Presentation(io.BytesIO(raw))
    tables = [s for s in deck.slides[0].shapes if getattr(s, "has_table", False)]
    assert tables, "pipe bullets must become a table"
    table = tables[0].table
    assert table.cell(0, 0).text == "项"
    assert table.cell(1, 1).text == "88"


def test_sandbox_helpers_make_cover_bullets_table() -> None:
    source = """
prs = Presentation()
add_title_slide(prs, "管理层决策会", "责任人 张三 · 2026-08-28")
add_content_slide(prs, "风险", ["汇率敞口", "客户流失", "库存积压"])
slide = add_content_slide(prs, "数字", [])
add_table(slide, [["项", "本季"], ["收入", "88"], ["同比", "-12%"]])
save_deck(prs)
"""
    raw = run_pptx_lib_source(source)
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) >= 3
    with zipfile.ZipFile(io.BytesIO(raw)) as zf:
        assert "ppt/presentation.xml" in zf.namelist()
    deck = Presentation(io.BytesIO(raw))
    title0 = deck.slides[0].shapes.title.text if deck.slides[0].shapes.title else ""
    assert "管理层决策会" in title0
    blob = "\n".join(
        (u.get("title") or "") + "\n" + "\n".join(u.get("bullets") or [])
        for u in (outline.get("units") or [])
        if isinstance(u, dict)
    )
    assert "张三" in blob
    tables = [s for s in deck.slides[2].shapes if getattr(s, "has_table", False)]
    assert tables
    assert tables[0].table.cell(1, 1).text == "88"


def test_sandbox_helpers_injected_no_import() -> None:
    src = (ROOT / "services/orchestrator/pico_orchestrator/office/sandbox_lib.py").read_text(
        encoding="utf-8"
    )
    assert "add_title_slide" in src
    assert "add_content_slide" in src
    assert "add_table" in src
    assert "RGBColor" in src
    ts = (ROOT / "services/true_pi_bridge/pico-gateway-tools.ts").read_text(encoding="utf-8")
    assert "add_title_slide" in ts
    assert "Prefer generate_pptx" not in ts
