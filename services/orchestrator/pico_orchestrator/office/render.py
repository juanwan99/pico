"""spec → legal OOXML bytes via python-docx / python-pptx / openpyxl."""

from __future__ import annotations

import io

from pico_orchestrator.office.comment import add_docx_comment
from pico_orchestrator.office.fill import fill_office_bytes
from pico_orchestrator.office.qa import verify_office_bytes
from pico_orchestrator.office.spec import OfficeSpec, Theme


def render_spec(spec: OfficeSpec, *, images: dict[str, bytes] | None = None) -> bytes:
    pictures = images or {}
    if spec.kind == "docx":
        raw = _render_docx(spec, pictures)
        ext = ".docx"
    elif spec.kind == "pptx":
        raw = _render_pptx(spec, pictures)
        ext = ".pptx"
    else:
        raw = _render_xlsx(spec)
        ext = ".xlsx"
    if spec.values:
        raw = fill_office_bytes(raw, ext, dict(spec.values))
    if spec.kind == "docx" and spec.comments:
        for item in spec.comments:
            raw = add_docx_comment(raw, paragraph_index=item.paragraph, text=item.text)
    check = verify_office_bytes(raw, ext)
    if not check["ok"]:
        raise ValueError(str(check.get("error") or "文档渲染失败。"))
    return raw


def _render_docx(spec: OfficeSpec, images: dict[str, bytes]) -> bytes:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.shared import Inches
    from docx.shared import RGBColor as DocxRGB

    doc = Document()
    _apply_docx_theme(doc, spec.theme)
    for block in spec.blocks:
        if block.type == "heading":
            level = 0 if block.level <= 0 else min(block.level, 4)
            doc.add_heading(block.text, level=level)
        elif block.type == "para":
            doc.add_paragraph(block.text)
        elif block.type == "table":
            table = doc.add_table(rows=len(block.rows), cols=len(block.rows[0]))
            table.style = "Table Grid"
            for r_i, row in enumerate(block.rows):
                for c_i, cell in enumerate(row):
                    table.rows[r_i].cells[c_i].text = cell
            _style_table_header(table, spec.theme)
        elif block.type == "image":
            raw = _need_image(images, block.artifact_id)
            doc.add_picture(io.BytesIO(raw), width=Inches(4.5))
            if block.text:
                doc.add_paragraph(block.text)
        elif block.type == "page_break":
            doc.add_page_break()
    if spec.theme and spec.theme.accent:
        color = _rgb(spec.theme.accent)
        if color is not None:
            for para in doc.paragraphs:
                style = str(getattr(getattr(para, "style", None), "name", "") or "")
                if style.lower().startswith("heading"):
                    for run in para.runs:
                        run.font.color.rgb = DocxRGB(*color)
    if spec.theme and spec.theme.heading_font:
        for para in doc.paragraphs:
            style = str(getattr(getattr(para, "style", None), "name", "") or "")
            if style.lower().startswith("heading"):
                for run in para.runs:
                    run.font.name = spec.theme.heading_font
                    run._element.rPr.rFonts.set(qn("w:eastAsia"), spec.theme.heading_font)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _render_pptx(spec: OfficeSpec, images: dict[str, bytes]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    deck = Presentation()
    layout = deck.slide_layouts[1] if len(deck.slide_layouts) > 1 else deck.slide_layouts[0]
    accent = _rgb(spec.theme.accent) if spec.theme else None
    body_font = spec.theme.body_font if spec.theme else None
    heading_font = spec.theme.heading_font if spec.theme else None
    for block in spec.blocks:
        if block.type != "slide":
            continue
        slide = deck.slides.add_slide(layout)
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            title_shape.text = block.title
            if heading_font or accent:
                for para in title_shape.text_frame.paragraphs:
                    for run in para.runs:
                        if heading_font:
                            run.font.name = heading_font
                            run.font.size = Pt(28)
                        if accent:
                            run.font.color.rgb = RGBColor(*accent)
        body = "\n".join(block.bullets)
        _set_pptx_body(
            slide,
            body,
            font_name=body_font,
            narrow_for_image=bool(block.image_artifact_id and block.bullets),
        )
        if block.image_artifact_id:
            raw = _need_image(images, block.image_artifact_id)
            _place_pptx_picture(slide, raw, has_bullets=bool(block.bullets))
    if not deck.slides:
        raise ValueError("PPT spec 没有可渲染的 slide。")
    buf = io.BytesIO()
    deck.save(buf)
    return buf.getvalue()


def _render_xlsx(spec: OfficeSpec) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    from openpyxl.utils import get_column_letter

    book = Workbook()
    default = book.active
    first = True
    accent = _rgb(spec.theme.accent) if spec.theme else None
    heading_font = spec.theme.heading_font if spec.theme else None
    body_font = spec.theme.body_font if spec.theme else None
    fill = None
    if accent is not None:
        fill = PatternFill("solid", fgColor=f"{accent[0]:02X}{accent[1]:02X}{accent[2]:02X}")
    header_font = Font(
        name=heading_font or "Calibri",
        bold=True,
        color="FFFFFF" if fill is not None else "000000",
    )
    data_font = Font(name=body_font or "Calibri") if body_font else None
    for block in spec.blocks:
        if block.type != "sheet":
            continue
        sheet = default if first else book.create_sheet()
        first = False
        sheet.title = (block.title or "Sheet1")[:31]
        row_i = 1
        if block.headers:
            for c_i, header in enumerate(block.headers, start=1):
                cell = sheet.cell(row=row_i, column=c_i, value=header)
                cell.font = header_font
                if fill is not None:
                    cell.fill = fill
            row_i += 1
        for row in block.rows:
            for c_i, raw in enumerate(row, start=1):
                cell = sheet.cell(row=row_i, column=c_i)
                _write_xlsx_cell(cell, raw)
                if data_font is not None:
                    cell.font = data_font
            row_i += 1
        if spec.marker and not any(spec.marker in "".join(row) for row in block.rows):
            sheet.cell(row=row_i, column=1, value=f"marker:{spec.marker}")
        width = max(len(block.headers), max((len(r) for r in block.rows), default=1), 1)
        for c_i in range(1, width + 1):
            sheet.column_dimensions[get_column_letter(c_i)].width = 14
    if first:
        raise ValueError("Excel spec 没有可渲染的 sheet。")
    buf = io.BytesIO()
    book.save(buf)
    return buf.getvalue()


def _write_xlsx_cell(cell: object, raw: str) -> None:
    text = "" if raw is None else str(raw)
    if text.startswith("="):
        cell.value = text
        return
    try:
        if text and text.lstrip("-").isdigit():
            cell.value = int(text)
            return
        if text and text.replace(".", "", 1).lstrip("-").isdigit() and text.count(".") == 1:
            cell.value = float(text)
            return
    except (TypeError, ValueError):
        pass
    cell.value = text


def _style_table_header(table: object, theme: Theme | None) -> None:
    if theme is None or not theme.accent or not getattr(table, "rows", None):
        return
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import RGBColor as DocxRGB

    color = _rgb(theme.accent)
    if color is None:
        return
    hex_color = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
    for cell in table.rows[0].cells:
        tc = cell._tc
        tc_pr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:fill"), hex_color)
        shd.set(qn("w:val"), "clear")
        tc_pr.append(shd)
        for para in cell.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = DocxRGB(255, 255, 255)
                if theme.heading_font:
                    run.font.name = theme.heading_font
                    run._element.rPr.rFonts.set(qn("w:eastAsia"), theme.heading_font)


# Default 10×7.5in Title-and-Content well. Pin all four edges when
# shrinking the body for a picture — setting only width/left collapses
# the placeholder to height=0 and the text paints over the title.
_PPTX_WELL_TOP_IN = 1.7
_PPTX_WELL_HEIGHT_IN = 5.2
_PPTX_BODY_LEFT_IN = 0.5
_PPTX_BODY_WIDTH_IN = 4.6
_PPTX_PIC_LEFT_IN = 5.15
_PPTX_PIC_WIDTH_IN = 4.5
_PPTX_HERO_LEFT_IN = 1.2
_PPTX_HERO_WIDTH_IN = 7.6
_PPTX_SLIDE_HEIGHT_IN = 7.5


def _place_pptx_picture(slide: object, raw: bytes, *, has_bullets: bool) -> None:
    """Put the picture in the content well. No extra spec fields."""
    from pptx.util import Inches

    if has_bullets:
        pic = slide.shapes.add_picture(
            io.BytesIO(raw),
            Inches(_PPTX_PIC_LEFT_IN),
            Inches(_PPTX_WELL_TOP_IN),
            width=Inches(_PPTX_PIC_WIDTH_IN),
        )
        _fit_pptx_picture(
            pic,
            left=Inches(_PPTX_PIC_LEFT_IN),
            top=Inches(_PPTX_WELL_TOP_IN),
            max_width=Inches(_PPTX_PIC_WIDTH_IN),
            max_height=Inches(_PPTX_WELL_HEIGHT_IN),
        )
        return
    pic = slide.shapes.add_picture(
        io.BytesIO(raw),
        Inches(_PPTX_HERO_LEFT_IN),
        Inches(_PPTX_WELL_TOP_IN),
        width=Inches(_PPTX_HERO_WIDTH_IN),
    )
    _fit_pptx_picture(
        pic,
        left=Inches(_PPTX_HERO_LEFT_IN),
        top=Inches(_PPTX_WELL_TOP_IN),
        max_width=Inches(_PPTX_HERO_WIDTH_IN),
        max_height=Inches(_PPTX_WELL_HEIGHT_IN),
    )


def _fit_pptx_picture(
    pic: object,
    *,
    left: int,
    top: int,
    max_width: int,
    max_height: int,
) -> None:
    pic.left = left
    pic.top = top
    width = int(getattr(pic, "width", 0) or 0)
    height = int(getattr(pic, "height", 0) or 0)
    if width > max_width and width > 0:
        height = int(height * (max_width / width))
        width = int(max_width)
    if height > max_height and height > 0:
        width = int(width * (max_height / height))
        height = int(max_height)
    pic.width = width
    pic.height = height


def _set_pptx_body(
    slide: object,
    body: str,
    *,
    font_name: str | None,
    narrow_for_image: bool = False,
) -> None:
    from pptx.util import Pt

    placeholders = getattr(slide, "placeholders", None)
    target = None
    if placeholders is not None:
        for shape in placeholders:
            idx = getattr(getattr(shape, "placeholder_format", None), "idx", None)
            if idx == 1 and getattr(shape, "has_text_frame", False):
                target = shape
                break
    if target is None:
        for shape in getattr(slide, "shapes", []):
            if shape is getattr(getattr(slide, "shapes", None), "title", None):
                continue
            if getattr(shape, "has_text_frame", False):
                target = shape
                break
    if target is None:
        return
    if narrow_for_image:
        from pptx.util import Inches

        target.left = Inches(_PPTX_BODY_LEFT_IN)
        target.top = Inches(_PPTX_WELL_TOP_IN)
        target.width = Inches(_PPTX_BODY_WIDTH_IN)
        target.height = Inches(_PPTX_WELL_HEIGHT_IN)
        if getattr(target, "text_frame", None) is not None:
            target.text_frame.word_wrap = True
    target.text = body
    if font_name:
        for para in target.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = font_name
                run.font.size = Pt(18)


def _apply_docx_theme(doc: object, theme: Theme | None) -> None:
    if theme is None or not theme.body_font:
        return
    from docx.oxml.ns import qn
    from docx.shared import Pt

    style = doc.styles["Normal"]
    style.font.name = theme.body_font
    style.font.size = Pt(12)
    style._element.rPr.rFonts.set(qn("w:eastAsia"), theme.body_font)


def _need_image(images: dict[str, bytes], artifact_id: str | None) -> bytes:
    if not artifact_id:
        raise ValueError("图缺少 artifact_id。")
    raw = images.get(artifact_id)
    if not raw:
        raise ValueError(f"找不到图片 {artifact_id}。请先 generate_image。")
    if raw[:8] != b"\x89PNG\r\n\x1a\n" and raw[:2] != b"\xff\xd8":
        raise ValueError("插入的图必须是 png 或 jpg。")
    return raw


def _rgb(value: str | None):
    if not value:
        return None
    text = value.strip().lstrip("#")
    if len(text) != 6:
        return None
    try:
        return tuple(int(text[i : i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return None
