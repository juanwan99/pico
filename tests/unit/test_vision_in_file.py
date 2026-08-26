"""T-VISION-IN-FILE: uploaded PPT/Word pictures become next-turn pixels.

Complex path (must pass): python-pptx embeds a real PNG → extract those
bytes → inspect_document remembers them → next /v1/chat/completions
``images=`` carries the same pixels. Counting pictures is not enough.
No /images/ fetch. No spec rebuild.
"""

from __future__ import annotations

import base64
import io
import os
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import pytest
from fastapi.testclient import TestClient

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "services" / "api"))
sys.path.insert(0, str(ROOT / "services" / "orchestrator"))

os.environ.setdefault("PICO_JWT_SECRET", "test-secret-at-least-32-bytes-long!!")

from app.main import app
from app.settings import Settings, get_settings
from pico_orchestrator.office.extract import extract_embedded_images
from pico_orchestrator.sandbox_s2 import PNG_MAGIC, encode_rgb_png
from pico_orchestrator.tools_builtin import build_default_gateway
from pico_orchestrator.usage_hook import bind_usage_context, reset_usage_context
from pico_orchestrator.vision import clear_conversation_images, conversation_images

FIXTURE_PNG = encode_rgb_png(
    48,
    32,
    bytes((220, 40, 40)) * (48 * 32),
    text_chunks={"in-file": "red"},
)
assert FIXTURE_PNG[:8] == PNG_MAGIC
FIXTURE_B64 = base64.b64encode(FIXTURE_PNG).decode("ascii")


@dataclass
class P:
    school_id: str
    membership_id: str
    scopes: list[str]


class MemoryArtifactStore:
    def __init__(self, *, run_id: str | None = None) -> None:
        self.rows: dict[tuple[str, str], list[dict[str, Any]]] = {}
        self._run_id = run_id
        self._task_id = "task-infile"

    def _rows(self, principal: P) -> list[dict[str, Any]]:
        return self.rows.setdefault((principal.school_id, principal.membership_id), [])

    async def write(
        self,
        principal: P,
        *,
        title: str,
        content: str | bytes,
        kind: str,
    ) -> dict[str, Any]:
        if isinstance(content, bytes):
            row = {
                "artifact_id": f"art-{sum(map(len, self.rows.values())) + 1}",
                "title": title,
                "content": None,
                "content_base64": base64.b64encode(content).decode("ascii"),
                "kind": kind,
                "run_id": self._run_id,
                "task_id": self._task_id,
                "size": len(content),
                "byte_size": len(content),
                "content_encoding": "base64",
            }
        else:
            row = {
                "artifact_id": f"art-{sum(map(len, self.rows.values())) + 1}",
                "title": title,
                "content": content,
                "kind": kind,
                "run_id": self._run_id,
                "task_id": self._task_id,
                "size": len(content.encode("utf-8")),
                "byte_size": len(content.encode("utf-8")),
                "content_encoding": "utf8",
            }
        self._rows(principal).append(row)
        return {k: v for k, v in row.items() if k not in {"content", "content_base64"}}

    async def read(
        self,
        principal: P,
        *,
        artifact_id: str | None = None,
        title: str | None = None,
    ) -> dict[str, Any] | None:
        for row in reversed(self._rows(principal)):
            if artifact_id and row["artifact_id"] == artifact_id:
                return row
            if not artifact_id and title and row["title"] == title:
                return row
        return None

    async def list(self, principal: P, *, limit: int) -> list[dict[str, Any]]:
        return [
            {k: v for k, v in row.items() if k not in {"content", "content_base64"}}
            for row in list(reversed(self._rows(principal)))[:limit]
        ]


def _pptx_with_picture() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "带图页"
    slide.shapes.add_picture(io.BytesIO(FIXTURE_PNG), Inches(1), Inches(1.5), width=Inches(5))
    buf = io.BytesIO()
    deck.save(buf)
    return buf.getvalue()


def _docx_with_picture() -> bytes:
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    doc.add_paragraph("正文旁边有图")
    doc.add_picture(io.BytesIO(FIXTURE_PNG), width=Inches(3))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _client() -> TestClient:
    settings = Settings(
        _env_file=None,
        pico_env="test",
        pico_openai_proxy_key="pico-dev",
        pico_allowed_models="pico-fast,pico-deep,pico-agent",
        pico_accept_test_issuer=True,
    )
    app.dependency_overrides[get_settings] = lambda: settings
    return TestClient(app)


@pytest.fixture(autouse=True)
def _clean_pending():
    clear_conversation_images()
    yield
    clear_conversation_images()
    app.dependency_overrides.pop(get_settings, None)


def test_extract_returns_the_same_pixels_not_just_a_count() -> None:
    raw = _pptx_with_picture()
    pictures = extract_embedded_images(raw, ".pptx")
    assert pictures, "extract found no picture"
    assert pictures[0] == FIXTURE_PNG
    word = extract_embedded_images(_docx_with_picture(), ".docx")
    assert word and word[0] == FIXTURE_PNG
    assert extract_embedded_images(raw, ".xlsx") == []


@pytest.mark.asyncio
async def test_inspect_uploaded_pptx_then_next_chat_sees_those_pixels(
    monkeypatch,
) -> None:
    store = MemoryArtifactStore(run_id="run-infile")
    gw = build_default_gateway(store)
    owner = P("school-a", "member-a", ["ai:run"])
    created = await store.write(
        owner, title="带图.pptx", content=_pptx_with_picture(), kind="pptx"
    )
    tok = bind_usage_context(
        school_id="school-a",
        membership_id="member-a",
        conversation_id="convo-infile",
    )
    try:
        outline = await gw.invoke(
            owner,
            "inspect_document",
            {"artifact_id": created["artifact_id"], "kind": "pptx"},
        )
        assert outline["images"] >= 1
        assert outline["extracted_images"] >= 1
        pending = conversation_images("convo-infile")
        assert pending, "inspect did not remember pixels"
        assert pending[0]["data"] == FIXTURE_B64
    finally:
        reset_usage_context(tok)

    captured: list[dict] = []

    async def fake_ledger(**kwargs):  # type: ignore[no-untyped-def]
        del kwargs
        return "task-in-1", "run-in-1"

    async def fake_run_and_collect(*_a, **kwargs):  # type: ignore[no-untyped-def]
        captured.append(kwargs)

        class R:
            status = "succeeded"
            final_text = "红块"
            error = None

        return R()

    async def fake_finalize(*_a, **_k):  # type: ignore[no-untyped-def]
        return None

    monkeypatch.setattr("app.openai_compat._ledger_task_run", fake_ledger)
    monkeypatch.setattr("app.openai_compat._run_and_collect", fake_run_and_collect)
    monkeypatch.setattr("app.openai_compat._finalize_run", fake_finalize)
    client = _client()
    try:
        r = client.post(
            "/v1/chat/completions",
            headers={
                "Authorization": "Bearer pico-dev",
                "X-Pico-Membership-Id": "test-member",
                "X-Conversation-Id": "convo-infile",
            },
            json={
                "model": "pico-fast",
                "messages": [{"role": "user", "content": "这页图是什么"}],
                "stream": False,
            },
        )
        assert r.status_code == 200, r.text
        images = (captured[0].get("images") or []) if captured else []
        assert images, "next chat never received extracted pixels"
        assert images[0]["data"] == FIXTURE_B64
        raw = base64.b64decode(images[0]["data"])
        assert raw == FIXTURE_PNG
    finally:
        app.dependency_overrides.pop(get_settings, None)
