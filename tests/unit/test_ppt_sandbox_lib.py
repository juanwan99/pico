"""T-PPT-SANDBOX-LIB: isolated python-pptx ceiling. Spec path stays default."""

from __future__ import annotations

import base64
import sys
import zipfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any
from unittest.mock import patch

import pytest

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "services" / "orchestrator"))

from pico_orchestrator.artifact_types import is_valid_ooxml_package
from pico_orchestrator.gateway import ToolError
from pico_orchestrator.office.inspect import inspect_office_bytes
from pico_orchestrator.office.pptx_helpers import ImagePathMap
from pico_orchestrator.office.sandbox_lib import assert_pptx_lib_source, run_pptx_lib_source
from pico_orchestrator.tools_builtin import build_default_gateway, openai_tool_schemas
from pico_orchestrator.true_pi.config import ALLOWED_GATEWAY_TOOLS
from pico_orchestrator.workbench_progress import (
    workbench_tool_result_line,
    workbench_tool_step_line,
)

ONE_PNG = (
    b"\x89PNG\r\n\x1a\n"
    b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
    b"\x00\x00\x00\x0cIDATx\x9cc``\x00\x00\x00\x04\x00\x01\xa3\x0a\x0d\xe4"
    b"\x00\x00\x00\x00IEND\xaeB`\x82"
)

THREE_SLIDE_WITH_IMAGE = """
prs = Presentation()
for title in ("封面", "配图", "结尾"):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
pic = None
for key in IMAGE_PATHS:
    pic = IMAGE_PATHS[key]
if pic:
    prs.slides[1].shapes.add_picture(pic, Inches(1), Inches(1.6), width=Inches(5))
save_deck(prs)
"""


@dataclass
class P:
    school_id: str = "school-a"
    membership_id: str = "member-a"
    scopes: list[str] | None = None

    def __post_init__(self) -> None:
        if self.scopes is None:
            self.scopes = ["ai:run"]


class MemoryArtifactStore:
    def __init__(self) -> None:
        self.rows: dict[tuple[str, str], list[dict[str, Any]]] = {}

    def _rows(self, principal: P) -> list[dict[str, Any]]:
        return self.rows.setdefault((principal.school_id, principal.membership_id), [])

    async def write(
        self,
        principal: P,
        *,
        title: str,
        content: str | bytes,
        kind: str,
    ) -> dict[str, Any]:
        if isinstance(content, bytes):
            row = {
                "artifact_id": f"art-{sum(map(len, self.rows.values())) + 1}",
                "title": title,
                "content": None,
                "content_base64": base64.b64encode(content).decode("ascii"),
                "kind": kind,
                "size": len(content),
                "byte_size": len(content),
                "content_encoding": "base64",
            }
        else:
            body = content
            row = {
                "artifact_id": f"art-{sum(map(len, self.rows.values())) + 1}",
                "title": title,
                "content": body,
                "kind": kind,
                "size": len(body.encode("utf-8")),
                "byte_size": len(body.encode("utf-8")),
                "content_encoding": "utf8",
            }
        self._rows(principal).append(row)
        return {k: v for k, v in row.items() if k not in {"content", "content_base64"}}

    async def read(
        self,
        principal: P,
        *,
        artifact_id: str | None,
        title: str | None,
    ) -> dict[str, Any] | None:
        for row in reversed(self._rows(principal)):
            if artifact_id and row["artifact_id"] == artifact_id:
                return row
            if not artifact_id and title and row["title"] == title:
                return row
        return None

    async def list(self, principal: P, *, limit: int) -> list[dict[str, Any]]:
        return [
            {k: v for k, v in row.items() if k != "content"}
            for row in list(reversed(self._rows(principal)))[:limit]
        ]


def _pptx_bytes(store: MemoryArtifactStore, owner: P, artifact_id: str) -> bytes:
    rows = store._rows(owner)
    row = next(r for r in rows if r["artifact_id"] == artifact_id)
    return base64.b64decode(row["content_base64"])


def test_allowlist_has_ceiling_not_bash() -> None:
    gw = build_default_gateway()
    names = set(gw.tools)
    assert "sandbox_pptx_lib" in names
    assert "generate_pptx_document" in names
    assert "sandbox_pptx_lib" in ALLOWED_GATEWAY_TOOLS
    assert "bash" not in ALLOWED_GATEWAY_TOOLS
    assert "bash" not in names
    schemas = {s["function"]["name"] for s in openai_tool_schemas(gw)}
    assert "sandbox_pptx_lib" in schemas
    assert len(ALLOWED_GATEWAY_TOOLS) == 27
    assert "generate_diagram" in ALLOWED_GATEWAY_TOOLS
    assert workbench_tool_step_line("sandbox_pptx_lib") == "正在沙箱排 PPT"
    assert workbench_tool_result_line("sandbox_pptx_lib", ok=True) == "已沙箱排 PPT"
    assert workbench_tool_result_line("sandbox_pptx_lib", ok=False) == "没沙箱排出 PPT"


def test_import_and_dunder_denied() -> None:
    with pytest.raises(ToolError) as ei:
        assert_pptx_lib_source("import os\nprs = Presentation()\nsave_deck(prs)")
    assert ei.value.code == "sandbox.exec_denied"
    with pytest.raises(ToolError) as ei:
        assert_pptx_lib_source("eval('1')")
    assert ei.value.code == "sandbox.exec_denied"
    with pytest.raises(ToolError) as ei:
        assert_pptx_lib_source("Presentation.__class__")
    assert ei.value.code == "sandbox.exec_denied"
    assert_pptx_lib_source(
        "from pptx import Presentation\nprs = Presentation()\nsave_deck(prs)"
    )


def test_blank_layout_go_title_body_writes() -> None:
    """Live F4 r2: blank + go() hit shapes.title None → sandbox.pptx_failed."""
    source = """
from pptx import Presentation, Inches, Pt
prs = Presentation()
blank = prs.slide_layouts[6]
def go(slide, title_txt, bullets):
    t = slide.shapes.title
    t.text = title_txt
    body = slide.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
slide = prs.slides.add_slide(blank)
box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.4))
box.text_frame.text = "办公简报：本周经营风险与下周动作"
slide = prs.slides.add_slide(blank)
go(slide, "本周经营风险总览", ["市场端询盘下降", "交付延迟风险"])
save_deck(prs)
"""
    raw = run_pptx_lib_source(source)
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) >= 2
    with zipfile.ZipFile(BytesIO(raw)) as zf:
        xml = b"".join(
            zf.read(name) for name in zf.namelist() if name.startswith("ppt/slides/")
        )
    blob = xml.decode("utf-8", errors="replace")
    assert "本周经营风险总览" in blob or "市场端询盘下降" in blob


def test_from_pptx_import_writes_real_deck() -> None:
    """Live F: document-skill `from pptx import Presentation` was exec_denied."""
    source = """
from pptx import Presentation
from pptx.util import Inches, Pt
prs = Presentation()
add_title_slide(prs, "本周经营风险与下周动作", "责任人：张三")
add_content_slide(prs, "风险总览", ["收入端延期", "毛利承压", "回款变慢"])
save_deck(prs)
"""
    raw = run_pptx_lib_source(source)
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) >= 2


def test_image_path_map_int_index() -> None:
    """Live F2: IMAGE_PATHS[0] was KeyError because the inject is a dict."""
    paths = ImagePathMap({"art-9": "/tmp/a.png", "art-2": "/tmp/b.png"})
    assert paths[0] == "/tmp/a.png"
    assert paths["art-9"] == "/tmp/a.png"
    assert paths[1] == "/tmp/b.png"
    assert 0 in paths
    assert paths.get(0) == "/tmp/a.png"
    empty = ImagePathMap()
    with pytest.raises(IndexError):
        empty[0]


def test_from_pptx_import_inches_pt_on_pptx_package() -> None:
    """Live F2: `from pptx import Presentation, Inches, Pt` was ImportError."""
    source = """
from pptx import Presentation, Inches, Pt
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "封面"
pic = IMAGE_PATHS[0]
slide.shapes.add_picture(pic, Inches(1), Inches(1.6), width=Inches(5))
save_deck(prs)
"""
    raw = run_pptx_lib_source(source, images={"art-cover": ONE_PNG})
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) >= 1
    assert int(outline["images"]) >= 1
    with zipfile.ZipFile(BytesIO(raw)) as zf:
        media = [name for name in zf.namelist() if name.startswith("ppt/media/")]
        assert media


def test_helper_aliases_title_image_and_table_prs() -> None:
    """Live F2: add_title_slide(image=) / add_table(prs=) / add_table(prs, rows)."""
    source = """
from pptx import Presentation, Inches, Pt
prs = Presentation()
add_title_slide(prs, "本周经营风险与下周动作", "责任人：张三", image=IMAGE_PATHS[0])
add_content_slide(prs, "风险总览", ["订单交付延期", "毛利率承压", "回款变慢"])
add_table(prs=prs, rows=[["项", "本周"], ["收入", "88"], ["毛利", "21"]])
add_table(prs, [["项", "上周"], ["回款", "12"]])
save_deck(prs)
"""
    raw = run_pptx_lib_source(source, images={"art-cover": ONE_PNG})
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) >= 4
    assert int(outline["images"]) >= 1
    from pptx import Presentation

    deck = Presentation(BytesIO(raw))
    tables = [s for sl in deck.slides for s in sl.shapes if getattr(s, "has_table", False)]
    assert len(tables) >= 2
    assert tables[0].table.cell(1, 1).text == "88"


def test_add_table_without_rows_still_typeerror() -> None:
    """Alias prs= does not invent a table when the grid is missing."""
    from pico_orchestrator.office.pptx_helpers import add_table
    from pptx import Presentation

    prs = Presentation()
    with pytest.raises(TypeError, match="rows"):
        add_table(prs=prs)


def test_live_placeholder_empty_shell_still_fails() -> None:
    source = (
        "from pptx import Presentation\n"
        "prs = Presentation()\n"
        "# 说明：此处使用沙箱其余逻辑由工具提供，本块只作占位\n"
        "save_deck(prs)"
    )
    with pytest.raises(ToolError) as ei:
        run_pptx_lib_source(source)
    assert ei.value.code == "sandbox.pptx_shell"


def test_empty_shell_fail_closed() -> None:
    with pytest.raises(ToolError) as ei:
        run_pptx_lib_source("prs = Presentation()\nsave_deck(prs)")
    assert ei.value.code == "sandbox.pptx_shell"
    assert "空壳" in ei.value.message or "没有可看" in ei.value.message
    with pytest.raises(ToolError) as ei:
        run_pptx_lib_source("prs = Presentation()")
    assert ei.value.code == "sandbox.pptx_empty"


@pytest.mark.asyncio
async def test_complex_image_then_sandbox_deck_has_media() -> None:
    """Complex task: generate_image → sandbox_pptx_lib → zip media + inspect."""
    store = MemoryArtifactStore()
    gw = build_default_gateway(store)
    owner = P()

    async def fake_image(_prompt: str) -> tuple[bytes, str]:
        return ONE_PNG, "png"

    with patch(
        "pico_orchestrator.tools_builtin.generate_image_bytes",
        fake_image,
    ):
        pictured = await gw.invoke(owner, "generate_image", {"prompt": "示意图"})
    aid = pictured["artifact_id"]
    assert aid

    out = await gw.invoke(
        owner,
        "sandbox_pptx_lib",
        {
            "source": THREE_SLIDE_WITH_IMAGE,
            "title": "上限.pptx",
            "image_artifact_ids": [aid],
        },
    )
    assert out["format"] == "pptx"
    assert out["via"] == "sandbox_pptx_lib"
    raw = _pptx_bytes(store, owner, out["artifact_id"])
    assert is_valid_ooxml_package(raw, ".pptx")
    with zipfile.ZipFile(BytesIO(raw)) as zf:
        media = [name for name in zf.namelist() if name.startswith("ppt/media/")]
        assert media, "sandbox deck must embed a picture, not a text-only shell"
        blob = zf.read(media[0])
        assert blob[:8] == b"\x89PNG\r\n\x1a\n"
        assert blob == ONE_PNG
    outline = inspect_office_bytes(raw, ".pptx")
    assert int(outline["slides"]) >= 3
    assert int(outline["images"]) >= 1


@pytest.mark.asyncio
async def test_generate_pptx_document_stays_default() -> None:
    gw = build_default_gateway(MemoryArtifactStore())
    owner = P()
    one = await gw.invoke(
        owner,
        "generate_pptx_document",
        {"title": "日常.pptx", "marker": "M-default", "body": "只有一页"},
    )
    assert one.get("format") == "pptx"
    assert one.get("via") != "sandbox_pptx_lib"


@pytest.mark.asyncio
async def test_tool_rejects_import_and_empty_shell() -> None:
    gw = build_default_gateway(MemoryArtifactStore())
    owner = P()
    with pytest.raises(ToolError) as ei:
        await gw.invoke(
            owner,
            "sandbox_pptx_lib",
            {"source": "import os\nprs = Presentation()\nsave_deck(prs)"},
        )
    assert ei.value.code == "sandbox.exec_denied"
    with pytest.raises(ToolError) as ei:
        await gw.invoke(
            owner,
            "sandbox_pptx_lib",
            {"source": "prs = Presentation()\nsave_deck(prs)"},
        )
    assert ei.value.code == "sandbox.pptx_shell"


def test_runner_is_python_not_host_bash() -> None:
    src = (ROOT / "services/orchestrator/pico_orchestrator/office/sandbox_lib.py").read_text(
        encoding="utf-8"
    )
    assert "sys.executable" in src
    assert "/bin/bash" not in src
    assert "shell=True" not in src
    assert "host bash" in src.lower() or "No host bash" in src or "禁止 host bash" in src
