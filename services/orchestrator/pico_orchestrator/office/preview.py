"""Content-box HTML projection of OOXML. Not a second Office OS.

Upstream: python-docx / python-pptx / openpyxl (same as render).
Codex-style: page/slide canvas only — no Writer/Impress chrome.
"""

from __future__ import annotations

import base64
import html
import io
from collections.abc import Callable, Iterator
from typing import Any

from pico_orchestrator.artifact_types import is_valid_ooxml_package
from pico_orchestrator.office.legacy import require_supported_office_ext

_MAX_SLIDES = 24
_MAX_IMAGES = 16
_MAX_IMAGE_BYTES = 1_500_000
_MAX_XLSX_ROWS = 40
_MAX_XLSX_COLS = 16

_OOXML_MISS = (AttributeError, ValueError, TypeError, KeyError, NotImplementedError)


def _opt(call: Callable[[], Any], default: Any = None) -> Any:
    """python-docx/pptx optional fields throw; skip that box, keep the page."""
    try:
        return call()
    except _OOXML_MISS:
        return default


def preview_office_html(raw: bytes, ext: str) -> str:
    suffix = require_supported_office_ext(ext)
    if suffix == ".pptx" or suffix == ".ppt":
        if not is_valid_ooxml_package(raw, ".pptx"):
            raise ValueError("不是真 PPT（OOXML）。")
        return _pptx_html(raw)
    if suffix == ".xlsx" or suffix == ".xls":
        return _xlsx_html(raw)
    if not is_valid_ooxml_package(raw, ".docx"):
        raise ValueError("不是真 Word（OOXML）。")
    return _docx_html(raw)


def _css() -> str:
    return """
:root{--ink:#1a1a1a;--muted:#6b6b6b;--line:#e4e4e4;--paper:#fff;--stage:#ececec}
*{box-sizing:border-box}
html,body{margin:0;background:var(--stage);color:var(--ink);
  font-family:-apple-system,BlinkMacSystemFont,"Segoe UI","PingFang SC","Microsoft YaHei",sans-serif}
.stage{padding:28px 18px 48px;display:flex;flex-direction:column;align-items:center;gap:22px}
.page{width:min(720px,100%);background:var(--paper);border:1px solid #ddd;
  border-radius:2px;padding:52px 56px 64px;box-shadow:0 8px 28px rgba(0,0,0,.06)}
.page h1{margin:0 0 14px;font-size:28px;font-weight:650;line-height:1.25}
.page h2{margin:18px 0 10px;font-size:20px;font-weight:650}
.page p{margin:0 0 10px;line-height:1.65;font-size:15px}
.page ul,.page ol{margin:0 0 12px;padding-left:1.25em;font-size:15px;line-height:1.65}
.page table{width:100%;border-collapse:collapse;margin:12px 0 18px;font-size:13px}
.page th,.page td{border:1px solid var(--line);padding:8px 10px;text-align:left}
.page th{background:#f4f4f4}
.page img{max-width:100%;height:auto;display:block;margin:12px 0}
.deck-card{width:min(960px,100%);display:flex;flex-direction:column;gap:8px}
.slide{width:100%;aspect-ratio:16/9;background:#f6f7f8;color:#18222b;
  border-radius:6px;overflow:hidden;position:relative;
  box-shadow:0 10px 30px rgba(0,0,0,.12);border:1px solid #d8d8d8}
.slide.dark{color:#f7fbff;border-color:transparent}
.slide .box{position:absolute;overflow:hidden;line-height:1.35;word-break:break-word}
.slide table{width:100%;height:100%;border-collapse:collapse;font-size:12px;background:#fff;color:#18222b}
.slide th,.slide td{border:1px solid #d5dee4;padding:5px 7px}
.slide th{background:#eef3f6;font-weight:650}
.slide img{width:100%;height:100%;object-fit:cover;display:block}
.kicker{font-size:11px;letter-spacing:.08em;color:var(--muted);margin:0}
"""


def _wrap(body: str, title: str) -> str:
    return (
        "<!doctype html><html lang='zh-CN'><head><meta charset='utf-8'>"
        f"<meta name='viewport' content='width=device-width, initial-scale=1'>"
        f"<title>{html.escape(title)}</title><style>{_css()}</style></head>"
        f"<body><div class='stage'>{body}</div></body></html>"
    )


def _data_uri(blob: bytes, content_type: str) -> str | None:
    if not blob or len(blob) > _MAX_IMAGE_BYTES:
        return None
    b64 = base64.b64encode(blob).decode("ascii")
    return f"data:{content_type};base64,{b64}"


def _guess_image_type(blob: bytes) -> str:
    if blob.startswith(b"\x89PNG"):
        return "image/png"
    if blob[:2] == b"\xff\xd8":
        return "image/jpeg"
    if blob[:4] == b"RIFF" and blob[8:12] == b"WEBP":
        return "image/webp"
    if blob[:6] in (b"GIF87a", b"GIF89a"):
        return "image/gif"
    return "image/jpeg"


def _docx_html(raw: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(raw))
    parts: list[str] = []
    image_n = 0
    for block in _docx_blocks(doc):
        if block[0] == "p":
            style, text, list_tag = block[1], block[2], block[3]
            if not text:
                continue
            if list_tag:
                parts.append(f"<{list_tag}><li>{html.escape(text)}</li></{list_tag}>")
                continue
            lowered = style.lower()
            if lowered.startswith("heading 1") or lowered == "title":
                tag = "h1"
            elif lowered.startswith("heading"):
                tag = "h2"
            else:
                tag = "p"
            parts.append(f"<{tag}>{html.escape(text)}</{tag}>")
        elif block[0] == "table":
            parts.append(_html_table(block[1]))
        elif block[0] == "image" and image_n < _MAX_IMAGES:
            uri = _data_uri(block[1], _guess_image_type(block[1]))
            if uri:
                image_n += 1
                parts.append(f"<img alt='' src='{uri}'>")
    if not parts:
        parts.append("<p class='kicker'>空文档</p>")
    return _wrap(f"<article class='page'>{''.join(parts)}</article>", "Word 内容框")


def _docx_list_tag(para: Any) -> str | None:
    try:
        from docx.oxml.ns import qn

        p_pr = para._element.find(qn("w:pPr"))
        if p_pr is None:
            return None
        num_pr = p_pr.find(qn("w:numPr"))
        if num_pr is None:
            return None
        return "ul"
    except _OOXML_MISS:
        return None


def _docx_blocks(doc: Any) -> list[tuple]:
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    out: list[tuple] = []
    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc)
            style = str(getattr(getattr(para, "style", None), "name", "") or "")
            text = str(para.text or "").strip()
            out.append(("p", style, text, _docx_list_tag(para)))
            for run in para.runs:
                try:
                    drawing = run._element.find(".//" + qn("a:blip"))
                except _OOXML_MISS:
                    drawing = None
                if drawing is None:
                    continue
                embed = drawing.get(qn("r:embed"))
                if not embed:
                    continue
                try:
                    part = doc.part.related_parts[embed]
                    blob = part.blob
                except _OOXML_MISS:
                    continue
                if blob:
                    out.append(("image", blob))
        elif child.tag == qn("w:tbl"):
            table = Table(child, doc)
            rows = [[(cell.text or "").strip() for cell in row.cells] for row in table.rows]
            out.append(("table", rows))
    return out


def _html_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    head, *body = rows
    th = "".join(f"<th>{html.escape(c)}</th>" for c in head)
    trs = [f"<tr>{th}</tr>"]
    for row in body:
        tds = "".join(f"<td>{html.escape(c)}</td>" for c in row)
        trs.append(f"<tr>{tds}</tr>")
    return f"<table>{''.join(trs)}</table>"


def _pptx_html(raw: bytes) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(raw))
    sw = float(prs.slide_width or 12192000)
    sh = float(prs.slide_height or 6858000)
    cards: list[str] = []
    image_n = 0
    for i, slide in enumerate(prs.slides, start=1):
        if i > _MAX_SLIDES:
            break
        bits: list[str] = []
        bg_css, tone = _slide_bg_css(slide)
        for shape, dx, dy in _walk_shapes(slide.shapes):
            box = _shape_box(shape, sw, sh, dx, dy)
            if box is None:
                continue
            style = (
                f"left:{box[0]:.2f}%;top:{box[1]:.2f}%;"
                f"width:{box[2]:.2f}%;height:{box[3]:.2f}%"
            )
            try:
                shape_type = shape.shape_type
            except _OOXML_MISS:
                shape_type = None
            if shape_type == MSO_SHAPE_TYPE.PICTURE and image_n < _MAX_IMAGES:
                try:
                    blob = shape.image.blob
                except _OOXML_MISS:
                    blob = b""
                uri = _data_uri(blob, _guess_image_type(blob)) if blob else None
                if uri:
                    image_n += 1
                    bits.append(f"<div class='box' style='{style}'><img alt='' src='{uri}'></div>")
                continue
            if getattr(shape, "has_table", False):
                rows = [
                    [(cell.text or "").strip() for cell in row.cells] for row in shape.table.rows
                ]
                bits.append(f"<div class='box' style='{style}'>{_html_table(rows)}</div>")
                continue
            inner = _shape_html(shape)
            if not inner:
                continue
            fallback_pt = max(12, min(36, int(box[3] * 1.6)))
            bits.append(
                f"<div class='box' style='{style};font-size:{fallback_pt}px'>{inner}</div>"
            )
        tone_class = " dark" if tone == "dark" else ""
        bg_attr = f" background:{bg_css};" if bg_css else ""
        cards.append(
            "<figure class='deck-card'>"
            f"<p class='kicker'>第 {i} 页</p>"
            f"<section class='slide{tone_class}' aria-label='第 {i} 页' "
            f"style='{bg_attr.strip()}'>{''.join(bits)}</section>"
            "</figure>"
        )
    if not cards:
        cards.append(
            "<figure class='deck-card'><p class='kicker'>第 1 页</p>"
            "<section class='slide' aria-label='第 1 页'>"
            "<div class='box' style='left:8%;top:40%'>空演示文稿</div></section></figure>"
        )
    return _wrap("".join(cards), "PPT 内容框")


def _walk_shapes(shapes: Any, dx: int = 0, dy: int = 0) -> Iterator[tuple[Any, int, int]]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        try:
            shape_type = shape.shape_type
        except _OOXML_MISS:
            shape_type = None
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                gdx = dx + int(shape.left or 0)
                gdy = dy + int(shape.top or 0)
                yield from _walk_shapes(shape.shapes, gdx, gdy)
            except _OOXML_MISS:
                continue
            continue
        yield shape, dx, dy


def _slide_bg_css(slide: Any) -> tuple[str, str]:
    try:
        from pptx.enum.dml import MSO_FILL

        fill = slide.background.fill
        if fill.type != MSO_FILL.SOLID:
            return "", "light"
        rgb = fill.fore_color.rgb
        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
        luma = (r * 299 + g * 587 + b * 114) / 1000
        return f"#{r:02x}{g:02x}{b:02x}", "dark" if luma < 140 else "light"
    except _OOXML_MISS:
        return "", "light"


def _shape_box(
    shape: Any, sw: float, sh: float, dx: int = 0, dy: int = 0
) -> tuple[float, float, float, float] | None:
    try:
        left = 100.0 * float((shape.left or 0) + dx) / sw
        top = 100.0 * float((shape.top or 0) + dy) / sh
        width = 100.0 * float(shape.width or 0) / sw
        height = 100.0 * float(shape.height or 0) / sh
    except _OOXML_MISS:
        return None
    if width <= 0.2 or height <= 0.2:
        return None
    return (left, top, width, height)


def _run_font_pt(run: Any) -> float | None:
    try:
        size = run.font.size
        if size is None:
            return None
        return float(size.pt)
    except _OOXML_MISS:
        return None


def _font_color_css(run: Any) -> str | None:
    try:
        rgb = run.font.color.rgb
        return f"#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}"
    except _OOXML_MISS:
        return None


def _shape_html(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    blocks: list[str] = []
    for para in shape.text_frame.paragraphs:
        bits: list[str] = []
        for run in para.runs:
            text = html.escape(str(run.text or ""))
            if not text:
                continue
            styles: list[str] = []
            pt = _run_font_pt(run)
            if pt:
                styles.append(f"font-size:{max(10, min(48, pt))}px")
            if _opt(lambda r=run: bool(r.font.bold)):
                styles.append("font-weight:700")
            color = _font_color_css(run)
            if color:
                styles.append(f"color:{color}")
            if styles:
                bits.append(f"<span style='{';'.join(styles)}'>{text}</span>")
            else:
                bits.append(text)
        if bits:
            blocks.append("".join(bits))
    if blocks:
        return "<br>".join(blocks)
    lines: list[str] = []
    for para in shape.text_frame.paragraphs:
        text = html.escape(str(para.text or "")).strip()
        if text:
            lines.append(text)
    return "<br>".join(lines)


def _xlsx_html(raw: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    try:
        sheet = wb.active
        rows: list[list[str]] = []
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i >= _MAX_XLSX_ROWS:
                break
            rows.append(["" if c is None else str(c) for c in row[:_MAX_XLSX_COLS]])
    finally:
        wb.close()
    body = _html_table(rows) if rows else "<p>空表</p>"
    return _wrap(f"<article class='page'><p class='kicker'>工作表</p>{body}</article>", "表格 内容框")
