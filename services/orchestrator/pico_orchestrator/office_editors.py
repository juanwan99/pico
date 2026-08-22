"""Thin PyPI adapters that edit existing OOXML. Not a self-built Office engine.

Upstream: python-docx / python-pptx. Load ledger bytes → mutate one field → save.
Forbidden: synthesizing a new package and calling it an edit of the original.
"""

from __future__ import annotations

import io

from pico_orchestrator.artifact_types import is_valid_ooxml_package


def _nonempty_paragraphs(document: object) -> list[object]:
    paras = getattr(document, "paragraphs", None) or []
    return [p for p in paras if str(getattr(p, "text", "") or "").strip()]


def edit_docx_bytes(
    raw: bytes,
    *,
    paragraph_index: int,
    text: str,
) -> bytes:
    """Replace one 1-based nonempty paragraph. Other paragraphs stay."""
    if not is_valid_ooxml_package(raw, ".docx"):
        raise ValueError("不是真 Word（OOXML）。请上传已有的 .docx。")
    new_text = (text or "").strip()
    if not new_text:
        raise ValueError("新段落不能为空。")
    if paragraph_index < 1:
        raise ValueError("段落序号从 1 起。")
    from docx import Document

    document = Document(io.BytesIO(raw))
    paras = _nonempty_paragraphs(document)
    if not paras:
        raise ValueError("这份 Word 没有可改的段落。")
    if paragraph_index > len(paras):
        raise ValueError(f"没有第 {paragraph_index} 段（共 {len(paras)} 段）。")
    paras[paragraph_index - 1].text = new_text
    out = io.BytesIO()
    document.save(out)
    data = out.getvalue()
    if not is_valid_ooxml_package(data, ".docx"):
        raise ValueError("改完后不是真 Word，未保存。")
    return data


def edit_pptx_title_bytes(
    raw: bytes,
    *,
    slide_index: int,
    new_title: str,
) -> bytes:
    """Replace the title of one 1-based slide. Other slides stay."""
    if not is_valid_ooxml_package(raw, ".pptx"):
        raise ValueError("不是真 PPT（OOXML）。请上传已有的 .pptx。")
    title = (new_title or "").strip()
    if not title:
        raise ValueError("新标题不能为空。")
    if slide_index < 1:
        raise ValueError("页码从 1 起。")
    from pptx import Presentation

    deck = Presentation(io.BytesIO(raw))
    slides = list(deck.slides)
    if not slides:
        raise ValueError("这份 PPT 没有可改的页。")
    if slide_index > len(slides):
        raise ValueError(f"没有第 {slide_index} 页（共 {len(slides)} 页）。")
    slide = slides[slide_index - 1]
    shape = getattr(slide.shapes, "title", None)
    if shape is not None and getattr(shape, "has_text_frame", False):
        shape.text = title
    else:
        replaced = False
        for item in slide.shapes:
            if getattr(item, "has_text_frame", False):
                item.text = title
                replaced = True
                break
        if not replaced:
            raise ValueError("该页没有标题可改。")
    out = io.BytesIO()
    deck.save(out)
    data = out.getvalue()
    if not is_valid_ooxml_package(data, ".pptx"):
        raise ValueError("改完后不是真 PPT，未保存。")
    return data
