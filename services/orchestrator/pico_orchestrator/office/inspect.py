"""Bytes → addressable outline. Model must not guess paragraph indexes.

School files are often irregular: two-row merged headers, mixed sheets,
Word/PPT nested tables, fill-blank prose with no table. Inspect reports
structure (merges, header band, a row window). It does not guess column
maps or dump the whole workbook into the model.
"""

from __future__ import annotations

import io
from typing import Any

from pico_orchestrator.artifact_types import is_valid_ooxml_package
from pico_orchestrator.office.comment import list_docx_comments
from pico_orchestrator.office.fill import leftover_placeholders
from pico_orchestrator.office.legacy import office_ext_for_bytes

_DEFAULT_HEADER_COLS = 64
_DEFAULT_PREVIEW_ROWS = 8
_MAX_HEADER_ROWS = 5
_MAX_PREVIEW_ROWS = 50
_MAX_COLS = 80
_MAX_MERGES = 40


def inspect_office_bytes(
    raw: bytes,
    ext: str,
    *,
    sheet: str | int | None = None,
    header_rows: int = 1,
    start_row: int | None = None,
    max_rows: int | None = None,
    max_cols: int | None = None,
) -> dict[str, object]:
    window = _normalize_window(
        sheet=sheet,
        header_rows=header_rows,
        start_row=start_row,
        max_rows=max_rows,
        max_cols=max_cols,
    )
    suffix = office_ext_for_bytes(ext, raw)
    if suffix == ".docx":
        return _inspect_docx(raw, window)
    if suffix == ".pptx":
        return _inspect_pptx(raw, window)
    return _inspect_xlsx(raw, window)


def _normalize_window(
    *,
    sheet: str | int | None,
    header_rows: int,
    start_row: int | None,
    max_rows: int | None,
    max_cols: int | None,
) -> dict[str, Any]:
    rows = 1 if header_rows is None else int(header_rows)
    if rows < 1 or rows > _MAX_HEADER_ROWS:
        raise ValueError(f"header_rows 必须是 1–{_MAX_HEADER_ROWS}。")
    preview_rows = _DEFAULT_PREVIEW_ROWS if max_rows is None else int(max_rows)
    if preview_rows < 1 or preview_rows > _MAX_PREVIEW_ROWS:
        raise ValueError(f"max_rows 必须是 1–{_MAX_PREVIEW_ROWS}。")
    cols = _DEFAULT_HEADER_COLS if max_cols is None else int(max_cols)
    if cols < 1 or cols > _MAX_COLS:
        raise ValueError(f"max_cols 必须是 1–{_MAX_COLS}。")
    start = None if start_row is None else int(start_row)
    if start is not None and start < 1:
        raise ValueError("start_row 必须从 1 起。")
    sheet_key: str | int | None
    if sheet is None or sheet == "":
        sheet_key = None
    elif isinstance(sheet, int) or (isinstance(sheet, str) and str(sheet).isdigit()):
        sheet_key = int(sheet)
        if sheet_key < 1:
            raise ValueError("sheet 序号必须从 1 起。")
    else:
        sheet_key = str(sheet).strip() or None
    return {
        "sheet": sheet_key,
        "header_rows": rows,
        "start_row": start,
        "max_rows": preview_rows,
        "max_cols": cols,
    }


def _inspect_docx(raw: bytes, window: dict[str, Any]) -> dict[str, object]:
    if not is_valid_ooxml_package(raw, ".docx"):
        raise ValueError("不是真 Word（OOXML）。")
    from docx import Document

    document = Document(io.BytesIO(raw))
    units: list[dict[str, object]] = []
    leftovers: list[str] = []
    para_n = 0
    table_n = 0
    image_n = 0
    max_cols = int(window["max_cols"])
    preview_rows = int(window["max_rows"])
    header_rows = int(window["header_rows"])
    for para in document.paragraphs:
        text = str(para.text or "").strip()
        if not text:
            continue
        para_n += 1
        leftovers.extend(leftover_placeholders(text))
        style = str(getattr(getattr(para, "style", None), "name", "") or "")
        kind = "heading" if style.lower().startswith("heading") else "para"
        units.append({"index": para_n, "kind": kind, "text": text[:240], "style": style})
    for table in document.tables:
        table_n += 1
        rows = [[(cell.text or "").strip() for cell in row.cells] for row in table.rows]
        for row in rows:
            for cell in row:
                leftovers.extend(leftover_placeholders(cell))
        units.append(_table_unit(table_n, rows, header_rows, preview_rows, max_cols))
    for rel in document.part.rels.values():
        reltype = str(getattr(rel, "reltype", "") or "")
        if "image" in reltype.lower():
            image_n += 1
            units.append({"index": image_n, "kind": "image"})
    comments = list_docx_comments(raw)
    units.extend(comments)
    return {
        "kind": "docx",
        "paragraphs": para_n,
        "tables": table_n,
        "images": image_n,
        "comments": len(comments),
        "placeholders": _unique(leftovers),
        "layout": "prose" if table_n == 0 else "tables",
        "units": units,
    }


def _inspect_pptx(raw: bytes, window: dict[str, Any]) -> dict[str, object]:
    if not is_valid_ooxml_package(raw, ".pptx"):
        raise ValueError("不是真 PPT（OOXML）。")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = Presentation(io.BytesIO(raw))
    units: list[dict[str, object]] = []
    leftovers: list[str] = []
    image_n = 0
    table_n = 0
    max_cols = int(window["max_cols"])
    preview_rows = int(window["max_rows"])
    header_rows = int(window["header_rows"])
    for i, slide in enumerate(deck.slides, start=1):
        title = ""
        shape = getattr(slide.shapes, "title", None)
        if shape is not None:
            title = str(getattr(shape, "text", "") or "").strip()
        leftovers.extend(leftover_placeholders(title))
        bullets: list[str] = []
        slide_images = 0
        max_w = 0
        max_h = 0
        slide_tables: list[dict[str, object]] = []
        for item in slide.shapes:
            if item is shape:
                continue
            if getattr(item, "has_table", False):
                table_n += 1
                grid = _pptx_table_rows(item.table)
                for row in grid:
                    for cell in row:
                        leftovers.extend(leftover_placeholders(cell))
                slide_tables.append(
                    _table_unit(table_n, grid, header_rows, preview_rows, max_cols)
                )
                continue
            if getattr(item, "has_text_frame", False):
                text = str(getattr(item, "text", "") or "").strip()
                if text:
                    leftovers.extend(leftover_placeholders(text))
                    bullets.extend(ln.strip() for ln in text.splitlines() if ln.strip())
            if getattr(item, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                slide_images += 1
                image_n += 1
                max_w = max(max_w, int(getattr(item, "width", 0) or 0))
                max_h = max(max_h, int(getattr(item, "height", 0) or 0))
        unit: dict[str, object] = {
            "index": i,
            "kind": "slide",
            "title": title[:240],
            "bullets": bullets[:12],
            "images": slide_images,
            "tables": len(slide_tables),
        }
        if slide_tables:
            unit["table_preview"] = slide_tables
        if slide_images:
            unit["image_width_emu"] = max_w
            unit["image_height_emu"] = max_h
        units.append(unit)
    return {
        "kind": "pptx",
        "slides": len(units),
        "images": image_n,
        "tables": table_n,
        "placeholders": _unique(leftovers),
        "units": units,
    }


def _inspect_xlsx(raw: bytes, window: dict[str, Any]) -> dict[str, object]:
    if not is_valid_ooxml_package(raw, ".xlsx"):
        raise ValueError("不是真 Excel（OOXML）。")
    from openpyxl import load_workbook

    book = load_workbook(io.BytesIO(raw), data_only=False, read_only=False)
    units: list[dict[str, object]] = []
    leftovers: list[str] = []
    formulas = 0
    wanted = window["sheet"]
    max_cols = int(window["max_cols"])
    header_rows = int(window["header_rows"])
    preview_rows = int(window["max_rows"])
    start_row = window["start_row"]
    for i, sheet in enumerate(book.worksheets, start=1):
        if wanted is not None and not _sheet_matches(wanted, i, sheet.title):
            continue
        merges = _merge_ranges(sheet)
        last_col = min(int(sheet.max_column or 0), max_cols)
        header_band = _header_band(sheet, header_rows, last_col)
        headers = _flatten_headers(header_band, merges=merges)
        blank_headers = sum(1 for h in headers if not str(h).strip())
        data_start = header_rows + 1
        window_start = data_start if start_row is None else max(int(start_row), data_start)
        sample: list[list[object]] = []
        sheet_formulas: list[str] = []
        formula_count = 0
        last_row = int(sheet.max_row or 0)
        window_end = min(last_row, window_start + preview_rows - 1) if last_row else 0
        for row_i, row in enumerate(sheet.iter_rows(values_only=False), start=1):
            cells = row[:last_col]
            for cell in cells:
                raw_val = cell.value
                if isinstance(raw_val, str):
                    leftovers.extend(leftover_placeholders(raw_val))
                    if raw_val.startswith("="):
                        formulas += 1
                        formula_count += 1
                        if len(sheet_formulas) < 12:
                            sheet_formulas.append(f"{cell.coordinate}:{raw_val}")
            if window_start <= row_i <= window_end:
                sample.append([_cell_preview(cell.value) for cell in cells])
        truncated = bool(last_row and window_end < last_row)
        leftover_rows = max(0, last_row - window_end) if truncated else 0
        units.append(
            {
                "index": i,
                "kind": "sheet",
                "name": sheet.title,
                "rows": last_row,
                "cols": int(sheet.max_column or 0),
                "headers": headers,
                "header_rows": header_band,
                "merges": merges,
                "preview": sample,
                "formulas": sheet_formulas,
                "formula_count": formula_count,
                "window": {
                    "start_row": window_start,
                    "end_row": window_end,
                    "max_cols": last_col,
                    "truncated": truncated,
                    "leftover_rows": leftover_rows,
                },
                "blank_header_cols": blank_headers,
                "irregular": bool(merges or header_rows > 1 or blank_headers),
            }
        )
    if wanted is not None and not units:
        raise ValueError("找不到指定的工作表。")
    return {
        "kind": "xlsx",
        "sheets": len(book.worksheets),
        "formulas": formulas,
        "placeholders": _unique(leftovers),
        "window": {
            "sheet": wanted,
            "header_rows": header_rows,
            "max_rows": preview_rows,
            "max_cols": max_cols,
        },
        "units": units,
    }


def _sheet_matches(wanted: str | int, index: int, title: str) -> bool:
    if isinstance(wanted, int):
        return wanted == index
    return str(wanted).strip() == str(title).strip()


def _merge_ranges(sheet: Any) -> list[str]:
    ranges = getattr(sheet, "merged_cells", None)
    if ranges is None:
        return []
    out: list[str] = []
    for item in list(ranges)[:_MAX_MERGES]:
        out.append(str(item))
    return out


def _header_band(sheet: Any, header_rows: int, max_cols: int) -> list[list[str]]:
    last_col = min(int(sheet.max_column or 0), max_cols)
    band: list[list[str]] = []
    for row_i in range(1, header_rows + 1):
        values: list[str] = []
        for col_i in range(1, last_col + 1):
            raw_val = sheet.cell(row_i, col_i).value
            values.append("" if raw_val is None else str(raw_val).strip()[:80])
        band.append(values)
    return band


def _flatten_headers(
    band: list[list[str]],
    *,
    merges: list[str] | None = None,
) -> list[str]:
    if not band:
        return []
    width = max(len(row) for row in band)
    filled = [list(row) + [""] * (width - len(row)) for row in band]
    for span in merges or []:
        _apply_merge_to_band(filled, span)
    headers: list[str] = []
    for col in range(width):
        parts = [row[col] for row in filled if row[col]]
        seen: list[str] = []
        for part in parts:
            if part not in seen:
                seen.append(part)
        headers.append(" / ".join(seen))
    return headers


def _apply_merge_to_band(band: list[list[str]], span: str) -> None:
    from openpyxl.utils.cell import range_boundaries

    try:
        min_col, min_row, max_col, max_row = range_boundaries(span)
    except (ValueError, TypeError):
        return
    if min_row < 1 or min_row > len(band):
        return
    width = len(band[0]) if band else 0
    if min_col < 1 or min_col > width:
        return
    value = ""
    for row_i in range(min_row, min(max_row, len(band)) + 1):
        for col_i in range(min_col, min(max_col, width) + 1):
            cell = band[row_i - 1][col_i - 1]
            if cell:
                value = cell
                break
        if value:
            break
    if not value:
        return
    for row_i in range(min_row, min(max_row, len(band)) + 1):
        for col_i in range(min_col, min(max_col, width) + 1):
            if not band[row_i - 1][col_i - 1]:
                band[row_i - 1][col_i - 1] = value


def _table_unit(
    index: int,
    rows: list[list[str]],
    header_rows: int,
    preview_rows: int,
    max_cols: int,
) -> dict[str, object]:
    width = 0
    for row in rows:
        width = max(width, len(row))
    clipped = [list(row[:max_cols]) + [""] * max(0, min(width, max_cols) - len(row)) for row in rows]
    band = clipped[:header_rows]
    headers = _flatten_headers(band)
    data_start = header_rows
    preview = clipped[data_start : data_start + preview_rows]
    leftover = max(0, len(clipped) - data_start - len(preview))
    return {
        "index": index,
        "kind": "table",
        "rows": len(rows),
        "cols": width,
        "headers": headers,
        "header_rows": band,
        "preview": preview,
        "window": {
            "start_row": data_start + 1,
            "end_row": data_start + len(preview),
            "truncated": leftover > 0,
            "leftover_rows": leftover,
        },
        "irregular": bool(header_rows > 1 or any(not h.strip() for h in headers)),
    }


def _pptx_table_rows(table: Any) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [str(getattr(cell, "text", "") or "").strip() for cell in row.cells]
        rows.append(cells)
    return rows


def _cell_preview(value: object) -> object:
    if value is None:
        return ""
    if isinstance(value, str):
        return value[:80]
    if isinstance(value, (int, float)):
        return value
    return str(value)[:80]


def _unique(items: list[str]) -> list[str]:
    seen: list[str] = []
    for item in items:
        if item not in seen:
            seen.append(item)
    return seen
